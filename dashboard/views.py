"""
views.py
--------
All views for the CHA Dashboard.
"""
import json
import csv
import io
from datetime import date

from django import forms as django_forms
from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse, JsonResponse
from django.contrib.auth import login, logout
from django.contrib.auth.decorators import login_required, user_passes_test
from django.contrib import messages
from django.db.models import Sum, Count, Q
from django.views.decorators.http import require_GET

from .models import UploadBatch, CHWRecord, SupervisionRecord
from .forms import LoginForm, UploadBatchForm
from .parsers import parse_chw_file, parse_supervision_file, compute_indicators


def healthz(request):
    from django.http import HttpResponse
    return HttpResponse("ok", status=200)


from django.http import JsonResponse, HttpResponse
from django.db.models import Sum, Count, Q
from django.views.decorators.http import require_GET

from .models import UploadBatch, CHWRecord, SupervisionRecord
from .forms import LoginForm, UploadBatchForm
from .parsers import parse_chw_file, parse_supervision_file, compute_indicators


# ---------------------------------------------------------------------------
# Auth
# ---------------------------------------------------------------------------

def login_view(request):
    if request.user.is_authenticated:
        return redirect('dashboard')
    form = LoginForm(request, data=request.POST or None)
    if request.method == 'POST' and form.is_valid():
        login(request, form.get_user())
        return redirect(request.GET.get('next', 'dashboard'))
    return render(request, 'dashboard/login.html', {'form': form})


def logout_view(request):
    logout(request)
    return redirect('login')


def is_uploader(user):
    return user.is_authenticated and (user.is_staff or user.is_superuser)


# ---------------------------------------------------------------------------
# Upload
# ---------------------------------------------------------------------------

@login_required
@user_passes_test(is_uploader)
def upload_view(request):
    form = UploadBatchForm(request.POST or None, request.FILES or None)
    if request.method == 'POST':
        if form.is_valid():
            batch = form.save(commit=False)
            batch.uploaded_by = request.user
            batch.year  = int(form.cleaned_data['year'])
            batch.month = int(form.cleaned_data['month'])
            batch.week_end_date = form.cleaned_data.get('week_end_date')
            batch.save()

            chw_rows, chw_errors = parse_chw_file(batch, request.FILES['chw_file'])
            sup_rows, sup_errors = parse_supervision_file(batch, request.FILES['supervision_file'])

            all_errors = chw_errors + sup_errors
            if all_errors:
                messages.warning(request, f"Upload completed with {len(all_errors)} row-level warnings. "
                                          f"CHW rows: {chw_rows}, Supervision rows: {sup_rows}.")
            else:
                messages.success(request, f"Upload successful! {chw_rows} CHW records and "
                                          f"{sup_rows} supervision records saved for {batch.label}.")
            return redirect('dashboard')
        else:
            messages.error(request, "Please correct the errors below.")

    batches = UploadBatch.objects.all().order_by('-year', '-month', '-week_start_date')
    return render(request, 'dashboard/upload.html', {'form': form, 'batches': batches})


@login_required
@user_passes_test(is_uploader)
def delete_batch_view(request, pk):
    batch = get_object_or_404(UploadBatch, pk=pk)
    if request.method == 'POST':
        label = batch.label
        batch.delete()
        messages.success(request, f"Batch '{label}' deleted successfully.")
    return redirect('upload')


# ---------------------------------------------------------------------------
# Dashboard
# ---------------------------------------------------------------------------

@login_required
def dashboard_view(request):
    # All available batches for the filter dropdowns
    batches = UploadBatch.objects.all().order_by('-year', '-month', '-week_start_date')

    # Selected filters
    selected_batch_id  = request.GET.get('batch')
    selected_county    = request.GET.get('county', '')
    selected_subcounty = request.GET.get('sub_county', '')
    selected_chu       = request.GET.get('chu', '')

    selected_batch = None
    indicators     = None
    filter_options = {}

    if selected_batch_id:
        selected_batch = get_object_or_404(UploadBatch, pk=selected_batch_id)

        chw_qs = CHWRecord.objects.filter(batch=selected_batch)
        sup_qs = SupervisionRecord.objects.filter(batch=selected_batch)

        # Cascading filter values for dropdowns
        filter_options['counties'] = (
            chw_qs.values_list('county', flat=True).distinct().order_by('county')
        )

        if selected_county:
            chw_qs = chw_qs.filter(county=selected_county)
            sup_qs = sup_qs.filter(county=selected_county)
            filter_options['sub_counties'] = (
                chw_qs.values_list('sub_county', flat=True).distinct().order_by('sub_county')
            )

        if selected_subcounty:
            chw_qs = chw_qs.filter(sub_county=selected_subcounty)
            sup_qs = sup_qs.filter(sub_county=selected_subcounty)
            filter_options['chus'] = (
                chw_qs.values_list('community_health_unit', flat=True).distinct().order_by('community_health_unit')
            )

        if selected_chu:
            chw_qs = chw_qs.filter(community_health_unit=selected_chu)
            sup_qs = sup_qs.filter(community_health_unit=selected_chu)

        indicators = compute_indicators(chw_qs, sup_qs, selected_batch.period_type)

    context = {
        'batches': batches,
        'selected_batch': selected_batch,
        'selected_batch_id': selected_batch_id or '',
        'selected_county': selected_county,
        'selected_subcounty': selected_subcounty,
        'selected_chu': selected_chu,
        'filter_options': filter_options,
        'indicators': indicators,
        'is_uploader': is_uploader(request.user),
    }
    return render(request, 'dashboard/dashboard.html', context)


# ---------------------------------------------------------------------------
# API endpoints (JSON) – called by JS for drill-down tables
# ---------------------------------------------------------------------------

@login_required
@require_GET
def api_inactive_chps(request):
    """Inactive CHPs for the selected batch."""
    batch_id   = request.GET.get('batch')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')

    if not batch_id:
        return JsonResponse({'error': 'batch required'}, status=400)

    qs = CHWRecord.objects.filter(batch_id=batch_id, is_active=False)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    data = list(qs.values(
        'county', 'sub_county', 'community_health_unit', 'chp_area', 'chw_name'
    ).order_by('county', 'sub_county', 'community_health_unit', 'chw_name'))

    return JsonResponse({'results': data, 'count': len(data)})


@login_required
@require_GET
def api_unsupervised(request):
    batch_id    = request.GET.get('batch')
    county      = request.GET.get('county', '')
    sub_county  = request.GET.get('sub_county', '')
    chu         = request.GET.get('chu', '')

    if not batch_id:
        return JsonResponse({'error': 'batch required'}, status=400)

    qs = CHWRecord.objects.filter(batch_id=batch_id, is_active=True, supervised=False)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    data = list(qs.values(
        'county', 'sub_county', 'community_health_unit', 'chp_area',
        'chw_name', 'chw_id', 'hh_visits', 'days_synced', 'supervision_visits'
    ).order_by('county', 'sub_county', 'community_health_unit', 'chw_name'))

    return JsonResponse({'results': data, 'count': len(data)})


@login_required
@require_GET
def api_low_performers(request):
    """Low performers split by supervision status."""
    batch_id    = request.GET.get('batch')
    county      = request.GET.get('county', '')
    sub_county  = request.GET.get('sub_county', '')
    chu         = request.GET.get('chu', '')
    group       = request.GET.get('group', 'unsupervised')  # 'unsupervised' or 'supervised'

    if not batch_id:
        return JsonResponse({'error': 'batch required'}, status=400)

    batch = get_object_or_404(UploadBatch, pk=batch_id)
    threshold = 50 if batch.period_type == 'monthly' else 12

    qs = CHWRecord.objects.filter(batch_id=batch_id, is_active=True, hh_visits__lt=threshold)
    if group == 'supervised':
        qs = qs.filter(supervised=True)
    else:
        qs = qs.filter(supervised=False)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    data = list(qs.values(
        'county', 'sub_county', 'community_health_unit', 'chp_area',
        'chw_name', 'chw_id', 'hh_visits', 'days_synced', 'supervision_visits'
    ).order_by('hh_visits', 'community_health_unit'))

    return JsonResponse({'results': data, 'count': len(data), 'threshold': threshold, 'group': group})


@login_required
@require_GET
def api_anc_gap(request):
    """CHPs with active pregnancies they did not fully visit."""
    from django.db.models import F as DjangoF
    batch_id    = request.GET.get('batch')
    county      = request.GET.get('county', '')
    sub_county  = request.GET.get('sub_county', '')
    chu         = request.GET.get('chu', '')

    if not batch_id:
        return JsonResponse({'error': 'batch required'}, status=400)

    qs = CHWRecord.objects.filter(batch_id=batch_id, is_active=True)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    # CHPs with active pregnancies but zero visits
    qs = qs.filter(active_pregnancies__gt=0, pregnancies_visited=0)

    data = list(qs.values(
        'county', 'sub_county', 'community_health_unit', 'chp_area',
        'chw_name', 'active_pregnancies', 'pregnancies_visited'
    ).order_by('community_health_unit', 'chw_name'))

    for row in data:
        row['gap'] = row['active_pregnancies']  # all unvisited since visited=0

    return JsonResponse({'results': data, 'count': len(data)})


@login_required
@require_GET
def api_supervised_3plus(request):
    """CHPs who received 3 or more supervision visits — data quality flag."""
    batch_id   = request.GET.get('batch')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')

    if not batch_id:
        return JsonResponse({'error': 'batch required'}, status=400)

    qs = CHWRecord.objects.filter(batch_id=batch_id, is_active=True, supervision_visits__gte=3)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    data = list(qs.values(
        'county', 'sub_county', 'community_health_unit', 'chp_area',
        'chw_name', 'hh_visits', 'supervision_visits'
    ).order_by('-supervision_visits', 'community_health_unit'))

    return JsonResponse({'results': data, 'count': len(data)})


@login_required
@require_GET
def api_u5_gap(request):
    """U5 assessment gap drill-downs."""
    batch_id   = request.GET.get('batch')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')
    gap_type   = request.GET.get('type', 'high_hh_low_u5')  # or 'high_u5_low_pos'

    if not batch_id:
        return JsonResponse({'error': 'batch required'}, status=400)

    qs = CHWRecord.objects.filter(
        batch_id=batch_id, is_active=True,
        registered_hhs__gt=0, registered_children_u5__gt=0
    )
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    rows = list(qs.values(
        'county', 'sub_county', 'community_health_unit', 'chp_area',
        'chw_name', 'registered_hhs', 'hh_visits',
        'registered_children_u5', 'num_u5_assessed', 'positive_diagnoses_u5', 'iccm_assessments'
    ))

    results = []
    for r in rows:
        hh_rate  = r['hh_visits'] / r['registered_hhs'] if r['registered_hhs'] else 0
        u5_rate  = r['num_u5_assessed'] / r['registered_children_u5'] if r['registered_children_u5'] else 0
        r['hh_rate_pct']  = round(hh_rate * 100, 1)
        r['u5_rate_pct']  = round(u5_rate * 100, 1)

        if gap_type == 'high_hh_low_u5' and hh_rate >= 0.7 and u5_rate < 0.4:
            results.append(r)
        elif gap_type == 'high_u5_low_pos' and u5_rate >= 0.8 and r['num_u5_assessed'] >= 10 and r['positive_diagnoses_u5'] == 0:
            results.append(r)

    results.sort(key=lambda x: x['u5_rate_pct'])
    return JsonResponse({'results': results, 'count': len(results)})


@login_required
@require_GET
def api_zero_pregnancies(request):
    """Active CHPs with zero active pregnancies registered."""
    batch_id   = request.GET.get('batch')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')

    if not batch_id:
        return JsonResponse({'error': 'batch required'}, status=400)

    qs = CHWRecord.objects.filter(batch_id=batch_id, is_active=True, active_pregnancies=0)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    data = list(qs.values(
        'county', 'sub_county', 'community_health_unit', 'chp_area',
        'chw_name', 'hh_visits', 'pregnancies_registered', 'active_pregnancies'
    ).order_by('community_health_unit', 'chw_name'))

    return JsonResponse({'results': data, 'count': len(data)})


@login_required
def download_zero_pregnancies(request):
    batch_id   = request.GET.get('batch')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')

    batch = get_object_or_404(UploadBatch, pk=batch_id)
    qs = CHWRecord.objects.filter(batch=batch, is_active=True, active_pregnancies=0)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="zero_active_pregnancies_{batch.label}.csv"'
    writer = csv.writer(response)
    writer.writerow(['County', 'Sub-County', 'Community Health Unit', 'CHP Area',
                     'CHP Name', 'HH Visits', 'Pregnancies Registered', 'Active Pregnancies'])
    for r in qs.order_by('community_health_unit', 'chw_name'):
        writer.writerow([r.county, r.sub_county, r.community_health_unit, r.chp_area,
                         r.chw_name, r.hh_visits, r.pregnancies_registered, r.active_pregnancies])
    return response


@login_required
@require_GET
def api_zero_positive(request):
    """Active CHPs with zero positive diagnoses U5."""
    batch_id   = request.GET.get('batch')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')

    if not batch_id:
        return JsonResponse({'error': 'batch required'}, status=400)

    qs = CHWRecord.objects.filter(batch_id=batch_id, is_active=True, positive_diagnoses_u5=0)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    data = list(qs.values(
        'county', 'sub_county', 'community_health_unit', 'chp_area',
        'chw_name', 'hh_visits', 'num_u5_assessed',
        'iccm_assessments', 'positive_diagnoses_u5'
    ).order_by('community_health_unit', 'chw_name'))

    return JsonResponse({'results': data, 'count': len(data)})


@login_required
def download_zero_positive(request):
    batch_id   = request.GET.get('batch')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')

    batch = get_object_or_404(UploadBatch, pk=batch_id)
    qs = CHWRecord.objects.filter(batch=batch, is_active=True, positive_diagnoses_u5=0)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="zero_positive_diagnoses_{batch.label}.csv"'
    writer = csv.writer(response)
    writer.writerow(['County', 'Sub-County', 'Community Health Unit', 'CHP Area',
                     'CHP Name', 'HH Visits', 'U5 Assessed', 'iCCM Assessments', 'Positive Diagnoses'])
    for r in qs.order_by('community_health_unit', 'chw_name'):
        writer.writerow([r.county, r.sub_county, r.community_health_unit, r.chp_area,
                         r.chw_name, r.hh_visits, r.num_u5_assessed,
                         r.iccm_assessments, r.positive_diagnoses_u5])
    return response


@login_required
@require_GET
def api_low_iccm(request):
    """Active CHPs with fewer than 5 iCCM assessments."""
    batch_id   = request.GET.get('batch')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')

    if not batch_id:
        return JsonResponse({'error': 'batch required'}, status=400)

    qs = CHWRecord.objects.filter(batch_id=batch_id, is_active=True, iccm_assessments__lt=5)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    data = list(qs.values(
        'county', 'sub_county', 'community_health_unit', 'chp_area',
        'chw_name', 'hh_visits', 'iccm_assessments',
        'registered_children_u5', 'num_u5_assessed'
    ).order_by('iccm_assessments', 'community_health_unit'))

    return JsonResponse({'results': data, 'count': len(data)})


@login_required
def download_low_iccm(request):
    batch_id   = request.GET.get('batch')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')

    batch = get_object_or_404(UploadBatch, pk=batch_id)
    qs = CHWRecord.objects.filter(batch=batch, is_active=True, iccm_assessments__lt=5)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="low_iccm_assessments_{batch.label}.csv"'
    writer = csv.writer(response)
    writer.writerow(['County', 'Sub-County', 'Community Health Unit', 'CHP Area',
                     'CHP Name', 'HH Visits', 'iCCM Assessments',
                     'Registered U5', 'U5 Assessed'])
    for r in qs.order_by('iccm_assessments', 'community_health_unit'):
        writer.writerow([r.county, r.sub_county, r.community_health_unit, r.chp_area,
                         r.chw_name, r.hh_visits, r.iccm_assessments,
                         r.registered_children_u5, r.num_u5_assessed])
    return response


@login_required
@require_GET
def api_same_day_flags(request):
    batch_id    = request.GET.get('batch')
    county      = request.GET.get('county', '')
    sub_county  = request.GET.get('sub_county', '')
    chu         = request.GET.get('chu', '')

    if not batch_id:
        return JsonResponse({'error': 'batch required'}, status=400)

    qs = SupervisionRecord.objects.filter(batch_id=batch_id)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    flags = (
        qs.values('county', 'sub_county', 'community_health_unit', 'visit_date')
        .annotate(count=Count('id'))
        .filter(count__gte=5)
        .order_by('-count', 'community_health_unit')
    )

    data = []
    for f in flags:
        data.append({
            'county': f['county'],
            'sub_county': f['sub_county'],
            'community_health_unit': f['community_health_unit'],
            'visit_date': f['visit_date'].strftime('%d %b %Y') if f['visit_date'] else '',
            'count': f['count'],
        })

    return JsonResponse({'results': data, 'count': len(data)})


# ---------------------------------------------------------------------------
# Download endpoints (CSV)
# ---------------------------------------------------------------------------

@login_required
def download_unsupervised(request):
    batch_id    = request.GET.get('batch')
    county      = request.GET.get('county', '')
    sub_county  = request.GET.get('sub_county', '')
    chu         = request.GET.get('chu', '')

    batch = get_object_or_404(UploadBatch, pk=batch_id)
    qs = CHWRecord.objects.filter(batch=batch, is_active=True, supervised=False)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="unsupervised_chps_{batch.label}.csv"'

    writer = csv.writer(response)
    writer.writerow(['County', 'Sub-County', 'Community Health Unit', 'CHP Area',
                     'CHP Name', 'CHP ID', 'HH Visits', 'Days Synced', 'Supervision Visits'])
    for r in qs.order_by('county', 'sub_county', 'community_health_unit'):
        writer.writerow([r.county, r.sub_county, r.community_health_unit, r.chp_area,
                         r.chw_name, r.chw_id, r.hh_visits, r.days_synced, r.supervision_visits])
    return response


@login_required
def download_low_performers(request):
    batch_id    = request.GET.get('batch')
    county      = request.GET.get('county', '')
    sub_county  = request.GET.get('sub_county', '')
    chu         = request.GET.get('chu', '')
    group       = request.GET.get('group', 'unsupervised')

    batch     = get_object_or_404(UploadBatch, pk=batch_id)
    threshold = 50 if batch.period_type == 'monthly' else 12
    qs = CHWRecord.objects.filter(batch=batch, is_active=True, hh_visits__lt=threshold)
    if group == 'supervised':
        qs = qs.filter(supervised=True)
    else:
        qs = qs.filter(supervised=False)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    label = 'low_performers_supervised' if group == 'supervised' else 'low_performers_not_supervised'
    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="{label}_{batch.label}.csv"'

    writer = csv.writer(response)
    writer.writerow(['County', 'Sub-County', 'Community Health Unit', 'CHP Area',
                     'CHP Name', 'CHP ID', 'HH Visits', 'Days Synced',
                     'Supervision Visits', f'Threshold (<{threshold})'])
    for r in qs.order_by('hh_visits', 'community_health_unit'):
        writer.writerow([r.county, r.sub_county, r.community_health_unit, r.chp_area,
                         r.chw_name, r.chw_id, r.hh_visits, r.days_synced,
                         r.supervision_visits, threshold])
    return response


@login_required
def download_anc_gap(request):
    from django.db.models import F as DjangoF
    batch_id    = request.GET.get('batch')
    county      = request.GET.get('county', '')
    sub_county  = request.GET.get('sub_county', '')
    chu         = request.GET.get('chu', '')

    batch = get_object_or_404(UploadBatch, pk=batch_id)
    qs = CHWRecord.objects.filter(batch=batch, is_active=True)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)
    qs = qs.filter(active_pregnancies__gt=0, pregnancies_visited=0
    ).order_by('community_health_unit', 'chw_name')

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="anc_gap_zero_visits_{batch.label}.csv"'

    writer = csv.writer(response)
    writer.writerow(['County', 'Sub-County', 'Community Health Unit', 'CHP Area',
                     'CHP Name', 'Active Pregnancies', 'Pregnancies Visited'])
    for r in qs:
        writer.writerow([r.county, r.sub_county, r.community_health_unit, r.chp_area,
                         r.chw_name, r.active_pregnancies, r.pregnancies_visited])
    return response


@login_required
def download_same_day_flags(request):
    batch_id    = request.GET.get('batch')
    county      = request.GET.get('county', '')
    sub_county  = request.GET.get('sub_county', '')
    chu         = request.GET.get('chu', '')

    batch = get_object_or_404(UploadBatch, pk=batch_id)
    qs = SupervisionRecord.objects.filter(batch=batch)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    flags = (
        qs.values('county', 'sub_county', 'community_health_unit', 'visit_date')
        .annotate(count=Count('id'))
        .filter(count__gte=5)
        .order_by('-count')
    )

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="same_day_flags_{batch.label}.csv"'

    writer = csv.writer(response)
    writer.writerow(['County', 'Sub-County', 'Community Health Unit', 'Visit Date', 'Supervisions Count'])
    for f in flags:
        writer.writerow([f['county'], f['sub_county'], f['community_health_unit'],
                         f['visit_date'], f['count']])
    return response


@login_required
def download_u5_gap(request):
    from django.db.models import F as DjangoF
    batch_id   = request.GET.get('batch')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')
    gap_type   = request.GET.get('type', 'high_hh_low_u5')

    batch = get_object_or_404(UploadBatch, pk=batch_id)
    qs = CHWRecord.objects.filter(
        batch=batch, is_active=True,
        registered_hhs__gt=0, registered_children_u5__gt=0
    )
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    rows = list(qs.values(
        'county', 'sub_county', 'community_health_unit', 'chp_area',
        'chw_name', 'registered_hhs', 'hh_visits',
        'registered_children_u5', 'num_u5_assessed', 'positive_diagnoses_u5'
    ))

    results = []
    for r in rows:
        hh_rate = r['hh_visits'] / r['registered_hhs'] if r['registered_hhs'] else 0
        u5_rate = r['num_u5_assessed'] / r['registered_children_u5'] if r['registered_children_u5'] else 0
        r['hh_rate_pct'] = round(hh_rate * 100, 1)
        r['u5_rate_pct'] = round(u5_rate * 100, 1)
        if gap_type == 'high_hh_low_u5' and hh_rate >= 0.7 and u5_rate < 0.4:
            results.append(r)
        elif gap_type == 'high_u5_low_pos' and u5_rate >= 0.8 and r['num_u5_assessed'] >= 10 and r['positive_diagnoses_u5'] == 0:
            results.append(r)

    if gap_type == 'high_hh_low_u5':
        filename = f'good_hh_low_u5_assessment_{batch.label}.csv'
        headers  = ['County', 'Sub-County', 'Community Health Unit', 'CHP Area', 'CHP Name',
                    'Registered HHs', 'HH Visits', 'HH Rate %',
                    'Registered U5', 'U5 Assessed', 'U5 Assessment Rate %']
    else:
        filename = f'high_u5_zero_positive_diagnoses_{batch.label}.csv'
        headers  = ['County', 'Sub-County', 'Community Health Unit', 'CHP Area', 'CHP Name',
                    'Registered U5', 'U5 Assessed', 'U5 Assessment Rate %', 'Positive Diagnoses']

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    writer = csv.writer(response)
    writer.writerow(headers)

    for r in sorted(results, key=lambda x: x['u5_rate_pct']):
        if gap_type == 'high_hh_low_u5':
            writer.writerow([r['county'], r['sub_county'], r['community_health_unit'], r['chp_area'],
                             r['chw_name'], r['registered_hhs'], r['hh_visits'], r['hh_rate_pct'],
                             r['registered_children_u5'], r['num_u5_assessed'], r['u5_rate_pct']])
        else:
            writer.writerow([r['county'], r['sub_county'], r['community_health_unit'], r['chp_area'],
                             r['chw_name'], r['registered_children_u5'], r['num_u5_assessed'],
                             r['u5_rate_pct'], r['positive_diagnoses_u5']])
    return response


@login_required
def download_supervised_3plus(request):
    batch_id   = request.GET.get('batch')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')

    batch = get_object_or_404(UploadBatch, pk=batch_id)
    qs = CHWRecord.objects.filter(batch=batch, is_active=True, supervision_visits__gte=3)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="supervised_3plus_{batch.label}.csv"'
    writer = csv.writer(response)
    writer.writerow(['County', 'Sub-County', 'Community Health Unit', 'CHP Area',
                     'CHP Name', 'HH Visits', 'Supervision Visits'])
    for r in qs.order_by('-supervision_visits', 'community_health_unit'):
        writer.writerow([r.county, r.sub_county, r.community_health_unit, r.chp_area,
                         r.chw_name, r.hh_visits, r.supervision_visits])
    return response


# ===========================================================================
# SYNC DASHBOARD VIEWS
# ===========================================================================

from .models import SyncUploadBatch, CHPSyncRecord
from .parsers import parse_sync_file, compute_sync_indicators

forms = django_forms  # alias so SyncUploadForm reads cleanly


class SyncUploadForm(django_forms.ModelForm):
    year  = django_forms.ChoiceField(choices=[(y, y) for y in range(2024, 2028)], initial=2026)
    month = django_forms.ChoiceField(choices=[(i, m) for i, m in [
        (1,'January'),(2,'February'),(3,'March'),(4,'April'),(5,'May'),(6,'June'),
        (7,'July'),(8,'August'),(9,'September'),(10,'October'),(11,'November'),(12,'December')
    ]])
    period_type = django_forms.ChoiceField(
        choices=[('monthly','Monthly'),('weekly','Weekly')],
        widget=django_forms.RadioSelect,
    )
    week_start_date = django_forms.DateField(
        required=False,
        widget=django_forms.DateInput(attrs={'type': 'date', 'class': 'form-input'}),
        label='Week Start Date',
        help_text='First day of the reporting period',
    )
    week_end_date = django_forms.DateField(
        required=False,
        widget=django_forms.DateInput(attrs={'type': 'date', 'class': 'form-input'}),
        label='Week End Date (optional)',
        help_text='Last day of the reporting period. Leave blank for standard week label.',
    )
    sync_file = django_forms.FileField(
        label='CHP Sync Report File (.xlsx)',
        widget=django_forms.FileInput(attrs={'accept': '.xlsx'}),
    )
    notes = django_forms.CharField(
        required=False,
        widget=django_forms.Textarea(attrs={'class': 'form-input', 'rows': 2}),
    )

    class Meta:
        model  = SyncUploadBatch
        fields = ['period_type', 'year', 'month', 'week_start_date', 'sync_file', 'notes']

    def clean(self):
        cleaned = super().clean()
        if cleaned.get('period_type') == 'weekly' and not cleaned.get('week_start_date'):
            self.add_error('week_start_date', 'Week start date is required for weekly uploads.')
        return cleaned


# Need to import forms here


@login_required
@user_passes_test(is_uploader)
def sync_upload_view(request):
    form = SyncUploadForm(request.POST or None, request.FILES or None)
    if request.method == 'POST' and form.is_valid():
        batch = form.save(commit=False)
        batch.uploaded_by = request.user
        batch.year  = int(form.cleaned_data['year'])
        batch.month = int(form.cleaned_data['month'])
        batch.week_end_date = form.cleaned_data.get('week_end_date')
        batch.save()

        rows, errors = parse_sync_file(batch, request.FILES['sync_file'])
        if errors:
            messages.warning(request, f"Upload completed with {len(errors)} warnings. {rows} rows saved.")
        else:
            messages.success(request, f"Sync report uploaded successfully! {rows} CHP records saved for {batch.label}.")
        return redirect('sync_dashboard')

    batches = SyncUploadBatch.objects.all().order_by('-year', '-month', '-week_start_date')
    return render(request, 'dashboard/sync_upload.html', {
        'form': form,
        'batches': batches,
        'is_uploader': True,
    })


@login_required
@user_passes_test(is_uploader)
def sync_delete_batch_view(request, pk):
    batch = get_object_or_404(SyncUploadBatch, pk=pk)
    if request.method == 'POST':
        label = batch.label
        batch.delete()
        messages.success(request, f"Sync batch '{label}' deleted.")
    return redirect('sync_upload')


def sync_dashboard_view(request):
    batches = SyncUploadBatch.objects.all()
    selected_batch_id  = request.GET.get('batch', '')
    selected_county    = request.GET.get('county', '')
    selected_subcounty = request.GET.get('sub_county', '')
    selected_chu       = request.GET.get('chu', '')

    selected_batch = None
    indicators     = None
    filter_options = {}

    if selected_batch_id:
        selected_batch = get_object_or_404(SyncUploadBatch, pk=selected_batch_id)
        qs = CHPSyncRecord.objects.filter(batch=selected_batch).exclude(
            county='').exclude(sub_county='').exclude(community_health_unit='')

        filter_options['counties'] = qs.values_list('county', flat=True).distinct().order_by('county')

        if selected_county:
            qs = qs.filter(county=selected_county)
            filter_options['sub_counties'] = qs.values_list('sub_county', flat=True).distinct().order_by('sub_county')

        if selected_subcounty:
            qs = qs.filter(sub_county=selected_subcounty)
            filter_options['chus'] = qs.values_list('community_health_unit', flat=True).distinct().order_by('community_health_unit')

        if selected_chu:
            qs = qs.filter(community_health_unit=selected_chu)

        indicators = compute_sync_indicators(qs)

    import json
    chu_data_json    = json.dumps(indicators['chus_sorted'] if indicators else [])
    county_data_json = json.dumps(indicators['counties']    if indicators else [])
    sc_data_json     = json.dumps(indicators['sub_counties'] if indicators else [])

    return render(request, 'dashboard/sync_dashboard.html', {
        'batches':            batches,
        'selected_batch':     selected_batch,
        'selected_batch_id':  selected_batch_id,
        'selected_county':    selected_county,
        'selected_subcounty': selected_subcounty,
        'selected_chu':       selected_chu,
        'filter_options':     filter_options,
        'indicators':         indicators,
        'chu_data_json':      chu_data_json,
        'county_data_json':   county_data_json,
        'sc_data_json':       sc_data_json,
        'show_county_table':  not selected_county,
        'show_sc_table':      bool(selected_county and not selected_subcounty),
        'show_chu_chart':     bool(selected_subcounty),
        'is_uploader':        is_uploader(request.user) if request.user.is_authenticated else False,
    })


@require_GET
def api_never_synced(request):
    batch_id    = request.GET.get('batch')
    county      = request.GET.get('county', '')
    sub_county  = request.GET.get('sub_county', '')
    chu         = request.GET.get('chu', '')

    if not batch_id:
        return JsonResponse({'error': 'batch required'}, status=400)

    qs = CHPSyncRecord.objects.filter(batch_id=batch_id, days_synced=0).exclude(
        county='').exclude(sub_county='').exclude(community_health_unit='')
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    data = list(qs.values(
        'county', 'sub_county', 'community_health_unit',
        'chp_name', 'username', 'reports_synced', 'last_sync_date'
    ).order_by('sub_county', 'community_health_unit', 'chp_name'))

    for row in data:
        row['last_sync_date'] = str(row['last_sync_date']) if row['last_sync_date'] else 'Never'

    return JsonResponse({'results': data, 'count': len(data)})


@login_required
def download_never_synced(request):
    batch_id    = request.GET.get('batch')
    county      = request.GET.get('county', '')
    sub_county  = request.GET.get('sub_county', '')
    chu         = request.GET.get('chu', '')

    batch = get_object_or_404(SyncUploadBatch, pk=batch_id)
    qs = CHPSyncRecord.objects.filter(batch=batch, days_synced=0).exclude(
        county='').exclude(sub_county='').exclude(community_health_unit='')
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="never_synced_{batch.label}.csv"'
    writer = csv.writer(response)
    writer.writerow(['County', 'Sub-County', 'Community Health Unit', 'CHP Name', 'Username', 'Days Synced', 'Reports Synced', 'Last Sync Date'])
    for r in qs.order_by('sub_county', 'community_health_unit', 'chp_name'):
        writer.writerow([r.county, r.sub_county, r.community_health_unit,
                         r.chp_name, r.username, r.days_synced, r.reports_synced,
                         r.last_sync_date or 'Never'])
    return response


@login_required
def download_chu_sync(request):
    batch_id    = request.GET.get('batch')
    county      = request.GET.get('county', '')
    sub_county  = request.GET.get('sub_county', '')
    chu         = request.GET.get('chu', '')

    from django.db.models import Count, Avg, Q

    batch = get_object_or_404(SyncUploadBatch, pk=batch_id)
    qs = CHPSyncRecord.objects.filter(batch=batch)
    if county:     qs = qs.filter(county=county)
    if sub_county: qs = qs.filter(sub_county=sub_county)
    if chu:        qs = qs.filter(community_health_unit=chu)

    chu_data = (
        qs.values('county', 'sub_county', 'community_health_unit')
        .annotate(
            total=Count('id'),
            synced=Count('id', filter=Q(days_synced__gte=1)),
            never=Count('id', filter=Q(days_synced=0)),
            avg_days=Avg('days_synced'),
            avg_reports=Avg('reports_synced'),
        )
        .order_by('sub_county', 'community_health_unit')
    )

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="chu_sync_{batch.label}.csv"'
    writer = csv.writer(response)
    writer.writerow(['County', 'Sub-County', 'Community Health Unit', 'Total CHPs',
                     'Synced', 'Never Synced', 'Sync Rate %', 'Avg Days Synced', 'Avg Reports Synced'])
    for c in chu_data:
        t = c['total'] or 1
        writer.writerow([
            c['county'], c['sub_county'], c['community_health_unit'],
            c['total'], c['synced'], c['never'],
            round(c['synced'] / t * 100, 1),
            round(c['avg_days'] or 0, 2),
            round(c['avg_reports'] or 0, 1),
        ])
    return response


# ===========================================================================
# SYNC COMPARISON VIEW
# ===========================================================================

def sync_compare_view(request):
    """Compare 2-4 sync batches with 5-category classification."""
    batches    = SyncUploadBatch.objects.all().order_by('-year', '-month', '-week_start_date')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')

    # Collect up to 4 batch IDs
    batch_ids = []
    for key in ['batch_a', 'batch_b', 'batch_c', 'batch_d']:
        val = request.GET.get(key, '').strip()
        if val and val not in batch_ids:
            batch_ids.append(val)

    comparison  = None
    filter_opts = {}
    selected_batch_ids = {
        'batch_a': request.GET.get('batch_a', ''),
        'batch_b': request.GET.get('batch_b', ''),
        'batch_c': request.GET.get('batch_c', ''),
        'batch_d': request.GET.get('batch_d', ''),
    }

    if len(batch_ids) >= 2:
        selected_batches = [get_object_or_404(SyncUploadBatch, pk=bid) for bid in batch_ids]
        n = len(selected_batches)

        def base_qs(batch):
            qs = CHPSyncRecord.objects.filter(batch=batch).exclude(
                county='').exclude(sub_county='').exclude(community_health_unit='')
            if county:     qs = qs.filter(county=county)
            if sub_county: qs = qs.filter(sub_county=sub_county)
            if chu:        qs = qs.filter(community_health_unit=chu)
            return qs

        def to_dict(qs):
            return {r['username']: r for r in qs.values(
                'username', 'chp_name', 'county', 'sub_county',
                'community_health_unit', 'days_synced', 'last_sync_date'
            )}

        dicts = [to_dict(base_qs(b)) for b in selected_batches]
        all_usernames = set()
        for d in dicts:
            all_usernames |= set(d.keys())

        # Labels for each batch column
        labels = [b.label for b in selected_batches]

        # 5 categories based on sync count across all batches
        synced_all   = []  # synced in every batch
        synced_most  = []  # synced in more than half but not all
        synced_some  = []  # synced in at least one but <= half
        never_synced = []  # never synced in any batch
        appeared_dropped = []  # only in one batch (new or dropped)

        for username in sorted(all_usernames):
            records = [d.get(username) for d in dicts]
            ref = next(r for r in records if r is not None)

            synced_flags = [bool(r and r['days_synced'] >= 1) for r in records]
            synced_count = sum(synced_flags)

            row = {
                'username':              username,
                'chp_name':              ref['chp_name'],
                'county':                ref['county'],
                'sub_county':            ref['sub_county'],
                'community_health_unit': ref['community_health_unit'],
                'synced_flags':          synced_flags,
                'synced_count':          synced_count,
                'days':  [r['days_synced'] if r else '—' for r in records],
                'dates': [str(r['last_sync_date']) if r and r['last_sync_date'] else 'Never' for r in records],
            }

            if synced_count == n:
                synced_all.append(row)
            elif synced_count == 0:
                never_synced.append(row)
            elif synced_count == 1:
                appeared_dropped.append(row)
            elif synced_count > n / 2:
                synced_most.append(row)
            else:
                synced_some.append(row)

        # Filter options
        filter_qs = CHPSyncRecord.objects.filter(batch=selected_batches[0]).exclude(
            county='').exclude(sub_county='').exclude(community_health_unit='')
        filter_opts['counties'] = filter_qs.values_list('county', flat=True).distinct().order_by('county')
        if county:
            filter_opts['sub_counties'] = filter_qs.filter(county=county).values_list('sub_county', flat=True).distinct().order_by('sub_county')
        if sub_county:
            filter_opts['chus'] = filter_qs.filter(county=county, sub_county=sub_county).values_list('community_health_unit', flat=True).distinct().order_by('community_health_unit')

        total = len(all_usernames)
        comparison = {
            'batches':          selected_batches,
            'labels':           labels,
            'n':                n,
            'total':            total,
            'synced_all':       synced_all,
            'synced_most':      synced_most,
            'synced_some':      synced_some,
            'never_synced':     never_synced,
            'appeared_dropped': appeared_dropped,
            'count_all':        len(synced_all),
            'count_most':       len(synced_most),
            'count_some':       len(synced_some),
            'count_never':      len(never_synced),
            'count_appeared':   len(appeared_dropped),
        }

    return render(request, 'dashboard/sync_compare.html', {
        'batches':          batches,
        'selected_batch_ids': selected_batch_ids,
        'comparison':       comparison,
        'filter_opts':      filter_opts,
        'selected_county':    county,
        'selected_subcounty': sub_county,
        'selected_chu':       chu,
        'is_uploader': is_uploader(request.user) if request.user.is_authenticated else False,
    })


@require_GET
def api_compare_download(request):
    """CSV download for any comparison category — supports 2-4 batches."""
    category   = request.GET.get('category', 'never_synced')
    county     = request.GET.get('county', '')
    sub_county = request.GET.get('sub_county', '')
    chu        = request.GET.get('chu', '')

    batch_ids = []
    for key in ['batch_a', 'batch_b', 'batch_c', 'batch_d']:
        val = request.GET.get(key, '').strip()
        if val and val not in batch_ids:
            batch_ids.append(val)

    if len(batch_ids) < 2:
        return HttpResponse('Need at least 2 batches', status=400)

    selected_batches = [get_object_or_404(SyncUploadBatch, pk=bid) for bid in batch_ids]
    n = len(selected_batches)

    def base_qs(batch):
        qs = CHPSyncRecord.objects.filter(batch=batch).exclude(
            county='').exclude(sub_county='').exclude(community_health_unit='')
        if county:     qs = qs.filter(county=county)
        if sub_county: qs = qs.filter(sub_county=sub_county)
        if chu:        qs = qs.filter(community_health_unit=chu)
        return qs

    dicts = [{r['username']: r for r in base_qs(b).values(
        'username', 'chp_name', 'county', 'sub_county',
        'community_health_unit', 'days_synced', 'last_sync_date'
    )} for b in selected_batches]

    all_usernames = set()
    for d in dicts:
        all_usernames |= set(d.keys())

    rows = []
    for username in sorted(all_usernames):
        records = [d.get(username) for d in dicts]
        ref = next(r for r in records if r is not None)
        synced_count = sum(1 for r in records if r and r['days_synced'] >= 1)

        include = False
        if category == 'synced_all'       and synced_count == n: include = True
        elif category == 'synced_most'    and synced_count > n/2 and synced_count < n: include = True
        elif category == 'synced_some'    and 0 < synced_count <= n/2 and synced_count > 1: include = True
        elif category == 'never_synced'   and synced_count == 0: include = True
        elif category == 'appeared_dropped' and synced_count == 1: include = True

        if include:
            rows.append({
                'username': username,
                'chp_name': ref['chp_name'],
                'county': ref['county'],
                'sub_county': ref['sub_county'],
                'community_health_unit': ref['community_health_unit'],
                'days': [r['days_synced'] if r else '—' for r in records],
                'dates': [str(r['last_sync_date']) if r and r['last_sync_date'] else 'Never' for r in records],
            })

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="sync_compare_{category}.csv"'
    writer = csv.writer(response)

    header = ['County', 'Sub-County', 'Community Health Unit', 'CHP Name', 'Username']
    for b in selected_batches:
        header += [f'Days Synced ({b.label})', f'Last Sync ({b.label})']
    writer.writerow(header)

    for r in rows:
        row = [r['county'], r['sub_county'], r['community_health_unit'], r['chp_name'], r['username']]
        for days, date in zip(r['days'], r['dates']):
            row += [days, date]
        writer.writerow(row)

    return response


# ===========================================================================
# WEEKLY PERFORMANCE SCORECARD
# ===========================================================================

# Hardcoded targets (universal across all counties/sub-counties/CHUs)
SCORECARD_TARGETS = {
    'active_chps':            {'target': None, 'unit': '',  'label': 'Active CHPs',              'higher_is_better': True, 'type': 'active_chps'},
    'hh_coverage_pct':        {'target': 85,   'unit': '%', 'label': 'HH Coverage',               'higher_is_better': True, 'type': 'simple'},
    'child_health':           {'target': None, 'unit': '',  'label': 'Child Health Indicators',   'higher_is_better': True, 'type': 'child_health'},
    'iccm_referrals':         {'target': None, 'unit': '',  'label': 'iCCM Referrals',            'higher_is_better': True, 'type': 'simple'},
    'iccm_referral_pct':      {'target': 90,   'unit': '%', 'label': 'iCCM Referrals Completed',  'higher_is_better': True, 'type': 'iccm_ref'},
    'pnc_ontime_pct':         {'target': 85,   'unit': '%', 'label': 'On Time PNC',               'higher_is_better': True, 'type': 'pnc'},
    'preg_registered_chp':    {'target': 1,    'unit': '',  'label': 'Preg Registered/CHP',       'higher_is_better': True, 'type': 'simple'},
    'sync_rate_pct':          {'target': 80,   'unit': '%', 'label': '% CHPs Syncing Weekly',     'higher_is_better': True, 'type': 'simple'},
    'supervision_pct':        {'target': 65,   'unit': '%', 'label': '% CHPs Supervised',         'higher_is_better': True, 'type': 'simple'},
}


def compute_scorecard_metrics(chw_qs, sync_qs=None):
    """Compute all scorecard metrics from a CHWRecord queryset."""
    from django.db.models import Sum, Count, Avg, Q

    active_qs   = chw_qs.filter(is_active=True)
    total_active = active_qs.count()
    total_all    = chw_qs.count()

    if total_all == 0:
        return None

    # 1. Active CHPs %
    active_pct = round(total_active / total_all * 100, 1) if total_all else 0

    # 2. HH Coverage % — avg across CHPs with registered HHs
    hh_qs = active_qs.filter(registered_hhs__gt=0)
    hh_metrics = hh_qs.aggregate(
        total_visits=Sum('hh_visits'),
        total_registered=Sum('registered_hhs'),
    )
    hh_coverage = round(
        (hh_metrics['total_visits'] or 0) / (hh_metrics['total_registered'] or 1) * 100, 1
    )

    # 3. Avg Positive Diagnoses per CHP
    pos_agg = active_qs.aggregate(total_pos=Sum('positive_diagnoses_u5'))
    avg_pos = round((pos_agg['total_pos'] or 0) / total_active, 2) if total_active else 0

    # 4. On Time PNC % = sum(PNC 48hr) / sum(Total Deliveries)
    #    Also PNC 3-7 day on-time
    pnc_agg = active_qs.filter(total_deliveries__gt=0).aggregate(
        pnc=Sum('pnc_48hr_ontime'),
        pnc_3_7d=Sum('pnc_3_7d_ontime'),
        del_total=Sum('total_deliveries'),
    )
    pnc_pct = round(
        (pnc_agg['pnc'] or 0) / (pnc_agg['del_total'] or 1) * 100, 1
    )
    pnc_3_7d_pct = round(
        (pnc_agg['pnc_3_7d'] or 0) / (pnc_agg['del_total'] or 1) * 100, 1
    )

    # 5. Pregnancies Registered per CHP
    preg_agg = active_qs.aggregate(total_preg=Sum('pregnancies_registered'))
    preg_per_chp = round((preg_agg['total_preg'] or 0) / total_active, 2) if total_active else 0

    # 6. % CHPs Supervised
    # Use whichever signal gives a non-100% result — 100% always indicates a data issue
    supervised_bool   = active_qs.filter(supervised=True).count()
    supervised_visits = active_qs.filter(supervision_visits__gt=0).count()

    if supervised_bool == total_active:
        # All CHPs show as supervised — likely a supervised-only extract bug
        # Fall back to supervision_visits count
        supervised = supervised_visits
    elif supervised_visits > 0 and supervised_visits != total_active:
        # supervision_visits has real data — use the higher of the two as they should agree
        supervised = max(supervised_bool, supervised_visits)
    else:
        supervised = supervised_bool

    sup_pct = round(supervised / total_active * 100, 1) if total_active else 0

    # 7. Sync rate — from sync queryset if provided
    sync_pct = None
    if sync_qs is not None:
        total_sync = sync_qs.count()
        synced     = sync_qs.filter(days_synced__gte=1).count()
        sync_pct   = round(synced / total_sync * 100, 1) if total_sync else 0

    # 8. Child Health Indicators
    child_agg = active_qs.aggregate(
        total_u5_assessed=Sum('num_u5_assessed'),
        total_registered_u5=Sum('registered_children_u5'),
        total_iccm=Sum('iccm_assessments'),
        total_pos=Sum('positive_diagnoses_u5'),
        total_referrals=Sum('iccm_referrals_total'),
        total_referrals_completed=Sum('iccm_referral_followup'),
        total_fever_cases=Sum('fever_cases'),
        total_fever_tested=Sum('fever_tested_rdt'),
    )
    total_u5_assessed    = child_agg['total_u5_assessed'] or 0
    total_registered_u5  = child_agg['total_registered_u5'] or 0
    total_iccm           = child_agg['total_iccm'] or 0
    total_referrals      = child_agg['total_referrals'] or 0
    total_ref_completed  = child_agg['total_referrals_completed'] or 0
    total_fever_cases    = child_agg['total_fever_cases'] or 0
    total_fever_tested   = child_agg['total_fever_tested'] or 0

    u5_assessment_pct = round(total_u5_assessed / total_registered_u5 * 100, 1) if total_registered_u5 else 0
    iccm_referral_pct = round(total_ref_completed / total_referrals * 100, 1) if total_referrals else None
    avg_pos_diag      = round((child_agg['total_pos'] or 0) / total_active, 1) if total_active else 0

    # PNC numerator/denominator for display
    pnc_numerator   = pnc_agg['pnc'] or 0
    pnc_denominator = pnc_agg['del_total'] or 0
    pnc_3_7d_numerator = pnc_agg['pnc_3_7d'] or 0

    # Active CHPs as % of total CHPs
    active_chps_pct = round(total_active / total_all * 100, 1) if total_all else 0

    return {
        'active_chps':         total_active,
        'total_chps':          total_all,
        'active_chps_pct':     active_chps_pct,
        'hh_coverage_pct':     hh_coverage,
        'hh_visits_total':     hh_metrics['total_visits'] or 0,
        'hh_registered_total': hh_metrics['total_registered'] or 0,
        'u5_assessment_pct':   u5_assessment_pct,
        'total_registered_u5': total_registered_u5,
        'total_u5_assessed':   total_u5_assessed,
        'iccm_assessments':    total_iccm,
        'avg_positive_diag':   avg_pos_diag,
        'total_positive_diag': child_agg['total_pos'] or 0,
        'fever_cases':         total_fever_cases,
        'fever_tested':        total_fever_tested,
        'iccm_referrals':      total_referrals,
        'iccm_referral_pct':   iccm_referral_pct,
        'iccm_ref_completed':  total_ref_completed,
        'pnc_ontime_pct':      pnc_pct,
        'pnc_numerator':       pnc_numerator,
        'pnc_denominator':     pnc_denominator,
        'pnc_3_7d_pct':        pnc_3_7d_pct,
        'pnc_3_7d_numerator':  pnc_3_7d_numerator,
        'preg_registered_chp': preg_per_chp,
        'supervision_pct':     sup_pct,
        'sync_rate_pct':       sync_pct,
    }


def get_colour(value, target, higher_is_better=True):
    """Return green/yellow/red based on % of target achieved. None target = no colour."""
    if value is None or target is None or target == 0:
        return 'grey'
    pct = value / target * 100 if higher_is_better else (2 * target - value) / target * 100
    if pct >= 100:   return 'green'
    if pct >= 50:    return 'yellow'
    return 'red'


def find_matching_sync_batch(chw_batch):
    """Find sync batch closest in time to a CHW batch."""
    if chw_batch is None:
        return None

    qs = SyncUploadBatch.objects.filter(
        year=chw_batch.year, month=chw_batch.month,
        period_type=chw_batch.period_type
    )

    if not qs.exists():
        # Try any sync batch from same year/month regardless of period type
        qs = SyncUploadBatch.objects.filter(
            year=chw_batch.year, month=chw_batch.month
        )

    if chw_batch.period_type == 'weekly' and chw_batch.week_start_date:
        # Prefer exact week_start_date match
        exact = qs.filter(week_start_date=chw_batch.week_start_date).first()
        if exact:
            return exact
        # Fall back to closest sync batch by week_start_date
        # Get all with a week_start_date and find the nearest
        dated = list(qs.exclude(week_start_date__isnull=True))
        if dated:
            return min(dated, key=lambda b: abs(
                (b.week_start_date - chw_batch.week_start_date).days
            ))

    return qs.order_by('-week_start_date', '-uploaded_at').first()


def auto_detect_batches(county=None, sub_county=None, chu=None):
    """
    Auto-detect which CHW batches to use for each scorecard column.
    Returns dict with keys: prev_month, weeks (list, oldest first)
    """
    all_batches = UploadBatch.objects.all().order_by('-year', '-month', '-week_start_date', '-uploaded_at')

    monthly  = list(all_batches.filter(period_type='monthly'))
    weekly   = list(all_batches.filter(period_type='weekly'))

    # Current month = month of the most recent weekly batch (or monthly if no weekly)
    current_month_ref = weekly[0] if weekly else (monthly[0] if monthly else None)

    weeks = []
    if current_month_ref:
        # All weekly batches in the same year/month as the most recent one, oldest first
        same_month_weekly = [b for b in weekly
                              if b.year == current_month_ref.year and b.month == current_month_ref.month]
        weeks = sorted(same_month_weekly, key=lambda b: b.week_start_date or date.min)

    # Previous month = most recent monthly batch from a strictly earlier month
    # OR most recent weekly batch from an earlier month if no monthly upload exists
    prev_month = None
    if current_month_ref:
        prev_monthly = [b for b in monthly if (b.year, b.month) < (current_month_ref.year, current_month_ref.month)]
        if prev_monthly:
            prev_month = prev_monthly[0]
        else:
            prev_weekly_other_month = [b for b in weekly if (b.year, b.month) < (current_month_ref.year, current_month_ref.month)]
            prev_month = prev_weekly_other_month[0] if prev_weekly_other_month else None

    return {
        'prev_month': prev_month,
        'weeks':      weeks,  # oldest to newest
    }


def auto_detect_monthly_batches():
    """
    Auto-detect up to 6 most recent monthly batches, oldest first.
    Falls back to including weekly batches if no monthly uploads exist for a month.
    """
    monthly = list(UploadBatch.objects.filter(
        period_type='monthly'
    ).order_by('-year', '-month')[:6])
    # Return oldest first
    return list(reversed(monthly))


def scorecard_view(request):
    """Weekly and Monthly performance scorecard view."""
    selected_counties    = request.GET.getlist('county')
    selected_subcounties = request.GET.getlist('sub_county')
    selected_chus        = request.GET.getlist('chu')

    # Toggle: weekly or monthly view
    scorecard_mode = request.GET.get('mode', 'weekly')

    # Manual override batch selectors
    override_prev_month   = request.GET.get('batch_prev_month', '')
    override_current_week = request.GET.get('batch_current', '')
    override_weeks        = request.GET.getlist('batch_week')
    override_months       = request.GET.getlist('batch_month')  # for monthly view

    all_batches         = UploadBatch.objects.all().order_by('-year', '-month', '-week_start_date')
    all_monthly_batches = UploadBatch.objects.filter(period_type='monthly').order_by('-year', '-month')

    # Auto-detect
    auto = auto_detect_batches()

    # Resolve previous month (single override still supported)
    batch_prev_month = UploadBatch.objects.filter(pk=override_prev_month).first() if override_prev_month else auto['prev_month']

    # Resolve weeks — if override_weeks provided, use those in order; else auto-detected weeks
    if override_weeks:
        batch_weeks = [UploadBatch.objects.filter(pk=wid).first() for wid in override_weeks]
        batch_weeks = [b for b in batch_weeks if b is not None]
    else:
        batch_weeks = auto['weeks']

    # Resolve monthly columns — up to 6 months oldest first
    if override_months:
        batch_months = [UploadBatch.objects.filter(pk=mid).first() for mid in override_months]
        batch_months = [b for b in batch_months if b is not None]
    else:
        batch_months = auto_detect_monthly_batches()

    def get_chw_qs(batch):
        if batch is None:
            return None
        qs = CHWRecord.objects.filter(batch=batch)
        if selected_counties:    qs = qs.filter(county__in=selected_counties)
        if selected_subcounties: qs = qs.filter(sub_county__in=selected_subcounties)
        if selected_chus:        qs = qs.filter(community_health_unit__in=selected_chus)
        return qs

    def get_sync_qs(chw_batch):
        sync_batch = find_matching_sync_batch(chw_batch)
        if sync_batch is None:
            return None
        qs = CHPSyncRecord.objects.filter(batch=sync_batch).exclude(
            county='').exclude(sub_county='').exclude(community_health_unit='')
        if selected_counties:    qs = qs.filter(county__in=selected_counties)
        if selected_subcounties: qs = qs.filter(sub_county__in=selected_subcounties)
        if selected_chus:        qs = qs.filter(community_health_unit__in=selected_chus)
        return qs

    metrics_prev_month = compute_scorecard_metrics(get_chw_qs(batch_prev_month), get_sync_qs(batch_prev_month)) if batch_prev_month else None

    # Compute metrics for each week column
    week_columns = []
    for b in batch_weeks:
        m = compute_scorecard_metrics(get_chw_qs(b), get_sync_qs(b))
        week_columns.append({'batch': b, 'metrics': m})

    # Compute metrics for each monthly column
    month_columns = []
    for b in batch_months:
        m = compute_scorecard_metrics(get_chw_qs(b), get_sync_qs(b))
        month_columns.append({'batch': b, 'metrics': m})

    # Current week = last week column; current month = last month column
    metrics_current_week  = week_columns[-1]['metrics']  if week_columns  else None
    batch_current_week    = week_columns[-1]['batch']    if week_columns  else None
    metrics_current_month = month_columns[-1]['metrics'] if month_columns else None
    batch_current_month   = month_columns[-1]['batch']   if month_columns else None

    # Build scorecard rows — both weekly and monthly use same indicators
    rows_weekly  = []
    rows_monthly = []

    for key, meta in SCORECARD_TARGETS.items():
        target   = meta['target']
        row_type = meta.get('type', 'simple')

        def make_cell(metrics, key=key, target=target, meta=meta, row_type=row_type):
            if metrics is None:
                return {'value': None, 'display': '—', 'colour': 'grey', 'pct_target': None, 'type': row_type}

            if row_type == 'active_chps':
                val   = metrics.get('active_chps', 0)
                pct   = metrics.get('active_chps_pct', 0)
                total = metrics.get('total_chps', 0)
                if pct >= 90:   colour = 'green'
                elif pct >= 70: colour = 'yellow'
                else:           colour = 'red'
                return {
                    'value': val, 'colour': colour, 'pct_target': round(pct, 1),
                    'pct_colour': colour,  # same threshold for % achieved column
                    'type': row_type,
                    'display': f"{val} ({pct}%)",
                    'detail': f"of {total} total CHPs",
                }

            elif row_type == 'child_health':
                ru5  = metrics.get('total_registered_u5', 0)
                au5  = metrics.get('total_u5_assessed', 0)
                upct = metrics.get('u5_assessment_pct', 0)
                iccm = metrics.get('iccm_assessments', 0)
                pos  = metrics.get('total_positive_diag', 0)
                avg  = metrics.get('avg_positive_diag', 0)
                fc   = metrics.get('fever_cases', 0)
                ft   = metrics.get('fever_tested', 0)
                colour = get_colour(upct, 100)

                # Two separate target achievements: U5 assessment % and avg sick children
                avg_target = 10
                avg_pct_target = round(avg / avg_target * 100, 1) if avg_target else None

                return {
                    'value': upct, 'colour': colour, 'pct_target': round(upct, 1),
                    'pct_target_2': avg_pct_target,
                    'pct_target_label':   'U5 Assessed',
                    'pct_target_2_label': 'Sick Children Avg',
                    'type': row_type,
                    'lines': [
                        ('U5 Pop', f"{ru5:,}"),
                        ('Assessed', f"{au5:,} ({upct}%)"),
                        ('iCCM Assessments', f"{iccm:,}"),
                        ('Sick Children (avg)', f"{pos:,} ({avg})"),
                        ('Fever', f"{fc:,} tested {ft:,}"),
                    ]
                }

            elif row_type == 'pnc':
                val48  = metrics.get('pnc_ontime_pct', 0)
                num48  = metrics.get('pnc_numerator', 0)
                den48  = metrics.get('pnc_denominator', 0)
                val37  = metrics.get('pnc_3_7d_pct', 0)
                num37  = metrics.get('pnc_3_7d_numerator', 0)
                colour = get_colour(val48, target)
                pct_target = round(val48 / target * 100, 1) if target else None
                return {
                    'value': val48, 'colour': colour, 'pct_target': pct_target, 'type': row_type,
                    'lines': [
                        ('PNC 48hr On-time', f"{val48}% ({num48}/{den48})"),
                        ('PNC 3-7d On-time', f"{val37}% ({num37}/{den48})"),
                    ]
                }

            elif row_type == 'iccm_ref':
                val  = metrics.get('iccm_referral_pct')
                comp = metrics.get('iccm_ref_completed', 0)
                tot  = metrics.get('iccm_referrals', 0)
                if val is None:
                    return {'value': None, 'display': '—', 'colour': 'grey', 'pct_target': None, 'type': row_type}
                pct_target = round(val / target * 100, 1) if target else None
                colour = get_colour(val, target)
                return {
                    'value': val, 'colour': colour, 'pct_target': pct_target, 'type': row_type,
                    'display': f"{comp} ({val}%)",
                }

            else:
                val = metrics.get(key)
                if val is None:
                    return {'value': None, 'display': '—', 'colour': 'grey', 'pct_target': None, 'type': row_type}
                unit = meta['unit']
                display = f"{val}{unit}" if unit == '%' else str(val)
                pct_target = round(val / target * 100, 1) if target else None
                colour = get_colour(val, target, meta['higher_is_better'])
                return {'value': val, 'display': display, 'colour': colour, 'pct_target': pct_target, 'type': row_type}

        # Monthly target display
        if row_type == 'child_health':
            target_display = f"100% of {(metrics_current_week or metrics_prev_month or {}).get('total_registered_u5', 0):,} U5s; avg {10} sick children"
        elif row_type == 'active_chps':
            total = (metrics_current_week or metrics_prev_month or {}).get('total_chps', 0)
            target_display = str(total) if total else '—'
        else:
            target_display = f"{target}{meta['unit']}" if target is not None else '—'

        base_row = {
            'key':    key,
            'label':  meta['label'],
            'target': target_display,
            'type':   row_type,
        }
        rows_weekly.append({
            **base_row,
            'prev_month': make_cell(metrics_prev_month),
            'weeks':      [make_cell(wc['metrics']) for wc in week_columns],
        })
        rows_monthly.append({
            **base_row,
            'months': [make_cell(mc['metrics']) for mc in month_columns],
        })

    # Filter options from the most data-rich batch
    filter_batch = batch_current_week or batch_prev_month
    filter_opts  = {}
    if filter_batch:
        fqs = CHWRecord.objects.filter(batch=filter_batch)
        filter_opts['counties'] = fqs.values_list('county', flat=True).distinct().order_by('county')
        if selected_counties:
            filter_opts['sub_counties'] = fqs.filter(county__in=selected_counties).values_list('sub_county', flat=True).distinct().order_by('sub_county')
        if selected_subcounties:
            filter_opts['chus'] = fqs.filter(sub_county__in=selected_subcounties).values_list('community_health_unit', flat=True).distinct().order_by('community_health_unit')

    # Inactive CHPs from the latest CHW batch (most recent week, else prev_month)
    latest_batch = batch_current_week or batch_prev_month
    inactive_chps = []
    if latest_batch:
        inactive_qs = CHWRecord.objects.filter(batch=latest_batch, is_active=False)
        if selected_counties:    inactive_qs = inactive_qs.filter(county__in=selected_counties)
        if selected_subcounties: inactive_qs = inactive_qs.filter(sub_county__in=selected_subcounties)
        if selected_chus:        inactive_qs = inactive_qs.filter(community_health_unit__in=selected_chus)

        # Try to get last sync date from matching sync batch
        sync_batch = find_matching_sync_batch(latest_batch)
        sync_lookup = {}
        if sync_batch:
            for rec in CHPSyncRecord.objects.filter(batch=sync_batch).values('username', 'last_sync_date'):
                if rec['username']:
                    sync_lookup[rec['username']] = rec['last_sync_date']

        for r in inactive_qs.order_by('county', 'sub_county', 'community_health_unit', 'chw_name'):
            last_sync = sync_lookup.get(r.username)
            inactive_chps.append({
                'chw_name':              r.chw_name,
                'county':                r.county,
                'sub_county':            r.sub_county,
                'community_health_unit': r.community_health_unit,
                'chp_area':              r.chp_area,
                'last_sync_date':        str(last_sync) if last_sync else '—',
            })

    return render(request, 'dashboard/scorecard.html', {
        'rows':                rows_weekly,
        'rows_monthly':        rows_monthly,
        'scorecard_mode':      scorecard_mode,
        'all_batches':         all_batches,
        'all_monthly_batches': all_monthly_batches,
        'batch_prev_month':    batch_prev_month,
        'week_batches':        [wc['batch'] for wc in week_columns],
        'month_batches':       [mc['batch'] for mc in month_columns],
        'batch_current_week':  batch_current_week,
        'batch_current_month': batch_current_month,
        'latest_batch':        batch_current_week or batch_current_month or batch_prev_month,
        'override_prev_month': override_prev_month,
        'override_weeks':      override_weeks,
        'override_months':     override_months,
        'filter_opts':         filter_opts,
        'selected_counties':   selected_counties,
        'selected_subcounties': selected_subcounties,
        'selected_chus':       selected_chus,
        'inactive_chps':       inactive_chps,
        'is_uploader': is_uploader(request.user) if request.user.is_authenticated else False,
    })


@login_required
def download_inactive_chps(request):
    counties     = request.GET.getlist('county')
    sub_counties = request.GET.getlist('sub_county')
    chus         = request.GET.getlist('chu')
    batch_id     = request.GET.get('batch')

    if batch_id:
        batch = get_object_or_404(UploadBatch, pk=batch_id)
    else:
        batch = UploadBatch.objects.order_by('-year', '-month', '-week_start_date').first()

    if not batch:
        return HttpResponse('No batch found', status=404)

    qs = CHWRecord.objects.filter(batch=batch, is_active=False)
    if counties:     qs = qs.filter(county__in=counties)
    if sub_counties: qs = qs.filter(sub_county__in=sub_counties)
    if chus:          qs = qs.filter(community_health_unit__in=chus)

    sync_batch  = find_matching_sync_batch(batch)
    sync_lookup = {}
    if sync_batch:
        for rec in CHPSyncRecord.objects.filter(batch=sync_batch).values('username', 'last_sync_date'):
            if rec['username']:
                sync_lookup[rec['username']] = rec['last_sync_date']

    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = f'attachment; filename="inactive_chps_{batch.label}.csv"'
    writer = csv.writer(response)
    writer.writerow(['CHP Name', 'County', 'Sub-County', 'Community Health Unit', 'CHP Area', 'Last Sync Date'])
    for r in qs.order_by('county', 'sub_county', 'community_health_unit', 'chw_name'):
        last_sync = sync_lookup.get(r.username, '—')
        writer.writerow([r.chw_name, r.county, r.sub_county, r.community_health_unit, r.chp_area, last_sync])
    return response


# ===========================================================================
# GAPS COMPARISON VIEW
# ===========================================================================

def get_gap_chps(batch, gap_key, county_list=None, sc_list=None, chu_list=None):
    """Return {username: row_data} for CHPs flagged for a given gap."""
    qs = CHWRecord.objects.filter(batch=batch, is_active=True)
    if county_list: qs = qs.filter(county__in=county_list)
    if sc_list:     qs = qs.filter(sub_county__in=sc_list)
    if chu_list:    qs = qs.filter(community_health_unit__in=chu_list)

    base_fields = ['username', 'chw_name', 'county', 'sub_county',
                   'community_health_unit', 'chp_area']

    if gap_key == 'not_supervised':
        qs = qs.filter(supervised=False)
    elif gap_key == 'lp_not_supervised':
        qs = qs.filter(hh_visits__lt=50, supervised=False)
    elif gap_key == 'lp_supervised':
        qs = qs.filter(hh_visits__lt=50, supervised=True)
    elif gap_key == 'supervised_3plus':
        qs = qs.filter(supervision_visits__gte=3)
    elif gap_key == 'good_hh_low_u5':
        # Use DB annotation for performance instead of Python loop
        from django.db.models import ExpressionWrapper, FloatField, F as DjF
        qs = qs.filter(registered_hhs__gt=0, registered_children_u5__gt=0)
        qs = qs.annotate(
            hh_rate=ExpressionWrapper(
                DjF('hh_visits') * 1.0 / DjF('registered_hhs'),
                output_field=FloatField()
            ),
            u5_rate=ExpressionWrapper(
                DjF('num_u5_assessed') * 1.0 / DjF('registered_children_u5'),
                output_field=FloatField()
            )
        ).filter(hh_rate__gte=0.7, u5_rate__lt=0.4)
    elif gap_key == 'high_u5_zero_diag':
        from django.db.models import ExpressionWrapper, FloatField, F as DjF
        qs = qs.filter(registered_children_u5__gt=0, num_u5_assessed__gte=10,
                       positive_diagnoses_u5=0)
        qs = qs.annotate(
            u5_rate=ExpressionWrapper(
                DjF('num_u5_assessed') * 1.0 / DjF('registered_children_u5'),
                output_field=FloatField()
            )
        ).filter(u5_rate__gte=0.8)
    elif gap_key == 'low_iccm':
        qs = qs.filter(iccm_assessments__lt=5)
    elif gap_key == 'zero_positive':
        qs = qs.filter(positive_diagnoses_u5=0)
    elif gap_key == 'anc_gap':
        qs = qs.filter(active_pregnancies__gt=0, pregnancies_visited=0)
    elif gap_key == 'zero_pregnancies':
        qs = qs.filter(active_pregnancies=0)
    else:
        return {}

    return {r['username']: r for r in qs.values(*base_fields) if r['username']}


def get_inactive_chps(batch, county_list=None, sc_list=None, chu_list=None):
    qs = CHWRecord.objects.filter(batch=batch, is_active=False)
    if county_list: qs = qs.filter(county__in=county_list)
    if sc_list:     qs = qs.filter(sub_county__in=sc_list)
    if chu_list:    qs = qs.filter(community_health_unit__in=chu_list)
    return {r['username']: r for r in qs.values(
        'username', 'chw_name', 'county', 'sub_county',
        'community_health_unit', 'chp_area') if r['username']}


def classify_chps(dicts, n):
    """
    Given list of {username: row} dicts (one per batch),
    classify each CHP as persistent/recurring/occasional.
    Returns dict: {username: {'count': N, 'row': {...}, 'periods': [T/F, ...]}}
    """
    all_usernames = set()
    for d in dicts:
        all_usernames |= set(d.keys())

    result = {'persistent': [], 'recurring': [], 'occasional': []}
    for username in sorted(all_usernames):
        periods = [username in d for d in dicts]
        count   = sum(periods)
        ref     = next((d[username] for d in dicts if username in d), None)
        if not ref:
            continue
        row = {**ref, 'flag_count': count, 'flag_total': n, 'periods': periods}
        if count == n:
            result['persistent'].append(row)
        elif count > n / 2:
            result['recurring'].append(row)
        else:
            result['occasional'].append(row)
    return result


GAP_INDICATORS = {
    'supervision': {
        'label': '🔍 Supervision',
        'indicators': [
            ('not_supervised',    'Not Supervised'),
            ('lp_not_supervised', 'Low Performers — Not Supervised'),
            ('lp_supervised',     'Low Performers — Supervised'),
            ('supervised_3plus',  'Supervised 3+ Times'),
        ]
    },
    'performance': {
        'label': '⚠️ Performance',
        'indicators': [
            ('good_hh_low_u5',   'Good HH Visits, Low U5 Assessment'),
            ('high_u5_zero_diag','High U5 Assessment, Zero Positive Diagnoses'),
            ('low_iccm',         'Low iCCM Assessments (<5)'),
            ('zero_positive',    'Zero Positive Diagnoses'),
        ]
    },
    'maternal': {
        'label': '🤱 Maternal Health',
        'indicators': [
            ('anc_gap',           'ANC Gap — Active Pregnancies, Zero Visits'),
            ('zero_pregnancies',  'Zero Active Pregnancies'),
        ]
    },
}


@login_required
def gaps_compare_view(request):
    all_batches  = UploadBatch.objects.all().order_by('-year', '-month', '-week_start_date')
    batch_ids    = [v for v in [request.GET.get('batch_a'), request.GET.get('batch_b'),
                                request.GET.get('batch_c')] if v]
    county_list  = request.GET.getlist('county')
    sc_list      = request.GET.getlist('sub_county')
    chu_list     = request.GET.getlist('chu')

    selected_batch_ids = {
        'batch_a': request.GET.get('batch_a', ''),
        'batch_b': request.GET.get('batch_b', ''),
        'batch_c': request.GET.get('batch_c', ''),
    }

    comparison   = None
    filter_opts  = {}

    if len(batch_ids) >= 2:
        try:
            batches = [get_object_or_404(UploadBatch, pk=bid) for bid in batch_ids]
            n       = len(batches)
            labels  = [b.label for b in batches]

            # Build filter options from first batch
            fqs = CHWRecord.objects.filter(batch=batches[0])
            filter_opts['counties'] = fqs.values_list('county', flat=True).distinct().order_by('county')
            if county_list:
                filter_opts['sub_counties'] = fqs.filter(county__in=county_list).values_list(
                    'sub_county', flat=True).distinct().order_by('sub_county')
            if sc_list:
                filter_opts['chus'] = fqs.filter(sub_county__in=sc_list).values_list(
                    'community_health_unit', flat=True).distinct().order_by('community_health_unit')

            kwargs = dict(county_list=county_list or None,
                          sc_list=sc_list or None,
                          chu_list=chu_list or None)

            # Compute all indicators
            categories = {}
            for cat_key, cat_meta in GAP_INDICATORS.items():
                indicators = {}
                for gap_key, gap_label in cat_meta['indicators']:
                    dicts = [get_gap_chps(b, gap_key, **kwargs) for b in batches]
                    classified = classify_chps(dicts, n)
                    indicators[gap_key] = {
                        'label': gap_label,
                        'persistent':  classified['persistent'],
                        'recurring':   classified['recurring'],
                        'occasional':  classified['occasional'],
                        'count_persistent': len(classified['persistent']),
                        'count_recurring':  len(classified['recurring']),
                        'count_occasional': len(classified['occasional']),
                    }
                categories[cat_key] = {
                    'label': cat_meta['label'],
                    'indicators': indicators,
                }

            # Inactive CHPs comparison
            inactive_dicts      = [get_inactive_chps(b, **kwargs) for b in batches]
            inactive_classified = classify_chps(inactive_dicts, n)
            inactive_ind = {
                'label':             'Inactive CHPs',
                'persistent':        inactive_classified['persistent'],
                'recurring':         inactive_classified['recurring'],
                'occasional':        inactive_classified['occasional'],
                'count_persistent':  len(inactive_classified['persistent']),
                'count_recurring':   len(inactive_classified['recurring']),
                'count_occasional':  len(inactive_classified['occasional']),
            }

            comparison = {
                'batches':   batches,
                'labels':    labels,
                'n':         n,
                'categories': categories,
                'inactive':  inactive_ind,
            }
        except Exception as e:
            import traceback
            comparison = {'error': str(e), 'traceback': traceback.format_exc()}

    return render(request, 'dashboard/gaps_compare.html', {
        'all_batches':       all_batches,
        'selected_batch_ids': selected_batch_ids,
        'comparison':        comparison,
        'filter_opts':       filter_opts,
        'selected_counties': county_list,
        'selected_subcounties': sc_list,
        'selected_chus':     chu_list,
        'is_uploader': is_uploader(request.user) if request.user.is_authenticated else False,
    })


@login_required
def gaps_compare_download(request):
    """Download comparison results as Excel with one sheet per indicator."""
    import openpyxl
    from openpyxl.styles import PatternFill, Font, Alignment
    from io import BytesIO

    batch_ids   = [v for v in [request.GET.get('batch_a'), request.GET.get('batch_b'),
                                request.GET.get('batch_c')] if v]
    county_list = request.GET.getlist('county') or None
    sc_list     = request.GET.getlist('sub_county') or None
    chu_list    = request.GET.getlist('chu') or None

    if len(batch_ids) < 2:
        return HttpResponse('Need at least 2 batches', status=400)

    batches = [get_object_or_404(UploadBatch, pk=bid) for bid in batch_ids]
    n       = len(batches)
    labels  = [b.label for b in batches]
    kwargs  = dict(county_list=county_list, sc_list=sc_list, chu_list=chu_list)

    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # remove default sheet

    header_fill   = PatternFill('solid', fgColor='1B3A6B')
    header_font   = Font(color='FFFFFF', bold=True)
    red_fill      = PatternFill('solid', fgColor='FEE2E2')
    yellow_fill   = PatternFill('solid', fgColor='FEF9C3')
    orange_fill   = PatternFill('solid', fgColor='FFEDD5')

    def write_sheet(ws, rows_by_group, labels):
        headers = ['Status', 'County', 'Sub-County', 'Community Health Unit',
                   'CHP Area', 'CHP Name', 'Username', 'Periods Flagged'] + \
                  [f'Flagged in {l}' for l in labels]
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=h)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal='center')

        row_num = 2
        fills   = {'persistent': red_fill, 'recurring': yellow_fill, 'occasional': orange_fill}
        labels_map = {'persistent': 'Persistent', 'recurring': 'Recurring', 'occasional': 'Occasional'}

        for group, rows in rows_by_group.items():
            for r in rows:
                periods_str = ', '.join(labels[i] for i, f in enumerate(r['periods']) if f)
                flag_cols   = ['Yes' if f else 'No' for f in r['periods']]
                data = [labels_map[group], r['county'], r['sub_county'],
                        r['community_health_unit'], r.get('chp_area', ''),
                        r['chw_name'], r['username'],
                        f"{r['flag_count']}/{r['flag_total']} — {periods_str}"] + flag_cols
                fill = fills[group]
                for col, val in enumerate(data, 1):
                    cell = ws.cell(row=row_num, column=col, value=val)
                    cell.fill = fill
                row_num += 1

        for col in ws.columns:
            ws.column_dimensions[col[0].column_letter].width = 18

    # Write indicator sheets
    for cat_key, cat_meta in GAP_INDICATORS.items():
        for gap_key, gap_label in cat_meta['indicators']:
            dicts      = [get_gap_chps(b, gap_key, **kwargs) for b in batches]
            classified = classify_chps(dicts, n)
            ws = wb.create_sheet(title=gap_label[:31])
            write_sheet(ws, classified, labels)

    # Inactive CHPs sheet
    inactive_dicts      = [get_inactive_chps(b, **kwargs) for b in batches]
    inactive_classified = classify_chps(inactive_dicts, n)
    ws = wb.create_sheet(title='Inactive CHPs')
    write_sheet(ws, inactive_classified, labels)

    buf = BytesIO()
    wb.save(buf)
    buf.seek(0)

    response = HttpResponse(
        buf.getvalue(),
        content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
    response['Content-Disposition'] = 'attachment; filename="gaps_comparison.xlsx"'
    return response


# ===========================================================================
# KPI REPORT UPLOAD VIEW
# ===========================================================================

@login_required
def kpi_upload_view(request):
    from .models import KPIReport
    from .forms import KPIReportForm
    from .kpi_parser import parse_kpi_report

    reports = KPIReport.objects.all().order_by('-report_year', '-report_month')
    success = error = None

    if request.method == 'POST':
        form = KPIReportForm(request.POST, request.FILES)
        if form.is_valid():
            report = form.save(commit=False)
            report.uploaded_by = request.user
            report.save()
            rows, errors = parse_kpi_report(report)
            if errors:
                error = f"Uploaded with warnings: {'; '.join(errors[:3])}"
            else:
                success = f"KPI Report uploaded — {rows} data points parsed."
        else:
            error = "Form invalid — check required fields."
    else:
        form = KPIReportForm()

    return render(request, 'dashboard/kpi_upload.html', {
        'form':    form,
        'reports': reports,
        'success': success,
        'error':   error,
        'is_uploader': is_uploader(request.user),
    })


@login_required
def kpi_delete_view(request, pk):
    from .models import KPIReport
    if request.method == 'POST':
        report = get_object_or_404(KPIReport, pk=pk)
        report.delete()
    return redirect('kpi_upload')


# ===========================================================================
# KPI SCORECARD DATA HELPERS
# ===========================================================================

MONTH_NAMES = {
    1:'Jan',2:'Feb',3:'Mar',4:'Apr',5:'May',6:'Jun',
    7:'Jul',8:'Aug',9:'Sep',10:'Oct',11:'Nov',12:'Dec'
}

def get_kpi_value(county, sub_county, metric_key, year, month, report_ids):
    """Fetch a single KPI value for a given geo/metric/period."""
    from .models import KPIDataPoint
    qs = KPIDataPoint.objects.filter(
        report_id__in=report_ids,
        county=county,
        sub_county=sub_county,
        metric_key=metric_key,
        year=year,
        month=month,
    ).first()
    return qs.value if qs else None


def compute_kpi_scorecard_metrics(county, sub_county, year, month, report_ids):
    """
    Compute scorecard-compatible metrics dict for a given geo/period
    from KPIDataPoint records.
    """
    def g(key):
        return get_kpi_value(county, sub_county, key, year, month, report_ids)

    active_chps     = g('active_chps')
    hh_cov          = g('hh_coverage_pct')
    u5_assessed     = g('u5_assessments')
    u5_total        = g('u5_children')
    sick_avg        = g('sick_children_avg')
    iccm_per        = g('iccm_assessments_per')
    iccm_ref_pct    = g('iccm_ref_completed')
    pnc_48          = g('pnc_48hr')
    pnc_37          = g('pnc_3_7d')
    deliveries      = g('facility_deliveries')
    preg_chp        = g('preg_per_chp')
    sync_pct        = g('sync_pct')
    sup_pct         = g('supervision_pct')

    # Compute derived
    u5_pct = round(u5_assessed / u5_total * 100, 1) if u5_total and u5_total > 0 else None
    hh_pct = round(hh_cov * 100, 1) if hh_cov is not None else None
    iccm_ref_pct_out = round(iccm_ref_pct * 100, 1) if iccm_ref_pct is not None else None
    sync_pct_out = round(sync_pct * 100, 1) if sync_pct is not None else None
    sup_pct_out  = round(sup_pct * 100, 1) if sup_pct is not None else None
    pnc_48_pct   = round(pnc_48 / deliveries * 100, 1) if pnc_48 and deliveries else None
    pnc_37_pct   = round(pnc_37 / deliveries * 100, 1) if pnc_37 and deliveries else None
    sick_avg_out = round(sick_avg, 2) if sick_avg is not None else None
    iccm_per_out = round(iccm_per, 1) if iccm_per is not None else None

    return {
        'active_chps':         int(active_chps) if active_chps is not None else None,
        'active_chps_pct':     None,  # no total CHPs in KPI file
        'total_chps':          None,
        'hh_coverage_pct':     hh_pct,
        'hh_visits_total':     None,
        'hh_registered_total': None,
        'total_registered_u5': int(u5_total) if u5_total else None,
        'total_u5_assessed':   int(u5_assessed) if u5_assessed else None,
        'u5_assessment_pct':   u5_pct,
        'iccm_assessments':    iccm_per_out,
        'avg_positive_diag':   sick_avg_out,
        'total_positive_diag': None,
        'fever_cases':         None,
        'fever_tested':        None,
        'iccm_referrals':      None,
        'iccm_referral_pct':   iccm_ref_pct_out,
        'iccm_ref_completed':  None,
        'pnc_ontime_pct':      pnc_48_pct,
        'pnc_numerator':       int(pnc_48) if pnc_48 else None,
        'pnc_denominator':     int(deliveries) if deliveries else None,
        'pnc_3_7d_pct':        pnc_37_pct,
        'pnc_3_7d_numerator':  int(pnc_37) if pnc_37 else None,
        'preg_registered_chp': round(preg_chp, 2) if preg_chp is not None else None,
        'sync_rate_pct':       sync_pct_out,
        'supervision_pct':     sup_pct_out,
    }


@login_required
def kpi_scorecard_view(request):
    """KPI Trends scorecard — reads from KPIDataPoint records."""
    from .models import KPIReport, KPIDataPoint

    # Filters
    selected_county    = request.GET.get('kpi_county', '')
    selected_subcounty = request.GET.get('kpi_subcounty', '')

    # Month selection — default last 6 available months
    selected_months = request.GET.getlist('kpi_month')  # list of "YYYY-MM"

    # Get all available reports
    all_reports = KPIReport.objects.all()
    report_ids  = list(all_reports.values_list('id', flat=True))

    # Get all available months from data
    available_months = list(
        KPIDataPoint.objects.filter(report_id__in=report_ids)
        .values('year', 'month')
        .distinct()
        .order_by('year', 'month')
    )

    # Build month columns — selected or last 6
    if selected_months:
        month_cols = []
        for ym in selected_months:
            try:
                y, m = ym.split('-')
                month_cols.append({'year': int(y), 'month': int(m),
                                   'label': f"{MONTH_NAMES.get(int(m), m)} {y}"})
            except Exception:
                pass
    else:
        last6 = available_months[-6:] if len(available_months) > 6 else available_months
        month_cols = [{'year': m['year'], 'month': m['month'],
                       'label': f"{MONTH_NAMES.get(m['month'], m['month'])} {m['year']}"}
                      for m in last6]

    # Geo resolution
    if selected_subcounty:
        county_key     = selected_county
        subcounty_key  = selected_subcounty
    elif selected_county:
        county_key    = selected_county
        subcounty_key = ''
    else:
        county_key    = ''
        subcounty_key = ''

    # Build filter options
    counties = list(
        KPIDataPoint.objects.filter(report_id__in=report_ids)
        .exclude(county='').values_list('county', flat=True)
        .distinct().order_by('county')
    )
    subcounties = []
    if selected_county:
        subcounties = list(
            KPIDataPoint.objects.filter(report_id__in=report_ids, county=selected_county)
            .exclude(sub_county='').values_list('sub_county', flat=True)
            .distinct().order_by('sub_county')
        )

    # Compute metrics per month column
    kpi_columns = []
    for mc in month_cols:
        m = compute_kpi_scorecard_metrics(
            county_key, subcounty_key, mc['year'], mc['month'], report_ids
        )
        kpi_columns.append({'label': mc['label'], 'year': mc['year'],
                            'month': mc['month'], 'metrics': m})

    # Build rows using same SCORECARD_TARGETS
    metrics_current = kpi_columns[-1]['metrics'] if kpi_columns else None
    rows = []

    for key, meta in SCORECARD_TARGETS.items():
        target   = meta['target']
        row_type = meta.get('type', 'simple')

        def make_kpi_cell(metrics, key=key, target=target, meta=meta, row_type=row_type):
            if metrics is None:
                return {'display': '—', 'colour': 'grey', 'pct_target': None, 'type': row_type}

            if row_type == 'active_chps':
                val = metrics.get('active_chps')
                if val is None:
                    return {'display': '—', 'colour': 'grey', 'pct_target': None, 'type': row_type}
                colour = 'grey'  # no total CHPs in KPI file
                return {'display': f"{val:,}", 'colour': colour, 'pct_target': None, 'type': row_type}

            elif row_type == 'child_health':
                ru5  = metrics.get('total_registered_u5')
                au5  = metrics.get('total_u5_assessed')
                upct = metrics.get('u5_assessment_pct')
                iccm = metrics.get('iccm_assessments')
                avg  = metrics.get('avg_positive_diag')
                colour = get_colour(upct, 100) if upct is not None else 'grey'
                avg_pct = round(avg / 10 * 100, 1) if avg is not None else None
                return {
                    'colour': colour, 'pct_target': upct,
                    'pct_target_2': avg_pct,
                    'pct_target_label': 'U5 Assessed',
                    'pct_target_2_label': 'Sick Children Avg',
                    'type': row_type,
                    'lines': [
                        ('U5 Pop',             f"{ru5:,}" if ru5 else '—'),
                        ('Assessed',           f"{au5:,} ({upct}%)" if au5 and upct else '—'),
                        ('iCCM Assess/CHP',    f"{iccm}" if iccm else '—'),
                        ('Sick Children (avg)', f"{avg}" if avg else '—'),
                        ('Fever',              '—'),
                    ]
                }

            elif row_type == 'pnc':
                val48 = metrics.get('pnc_ontime_pct')
                num48 = metrics.get('pnc_numerator')
                den   = metrics.get('pnc_denominator')
                val37 = metrics.get('pnc_3_7d_pct')
                num37 = metrics.get('pnc_3_7d_numerator')
                colour = get_colour(val48, target) if val48 is not None else 'grey'
                pct_target = round(val48 / target * 100, 1) if val48 and target else None
                return {
                    'colour': colour, 'pct_target': pct_target, 'type': row_type,
                    'lines': [
                        ('PNC 48hr On-time', f"{val48}% ({num48}/{den})" if val48 else '—'),
                        ('PNC 3-7d On-time', f"{val37}% ({num37}/{den})" if val37 else '—'),
                    ]
                }

            elif row_type == 'iccm_ref':
                val = metrics.get('iccm_referral_pct')
                if val is None:
                    return {'display': '—', 'colour': 'grey', 'pct_target': None, 'type': row_type}
                pct_target = round(val / target * 100, 1) if target else None
                colour = get_colour(val, target)
                return {'display': f"{val}%", 'colour': colour, 'pct_target': pct_target, 'type': row_type}

            else:
                val = metrics.get(key)
                if val is None:
                    return {'display': '—', 'colour': 'grey', 'pct_target': None, 'type': row_type}
                unit = meta['unit']
                display = f"{val}{unit}" if unit == '%' else str(val)
                pct_target = round(val / target * 100, 1) if target else None
                colour = get_colour(val, target, meta.get('higher_is_better', True))
                return {'display': display, 'colour': colour, 'pct_target': pct_target, 'type': row_type}

        # Target display
        if row_type == 'child_health':
            ru5 = (metrics_current or {}).get('total_registered_u5', 0)
            target_display = f"100% of {ru5:,} U5s; avg 10 sick/CHP" if ru5 else "—"
        elif row_type == 'active_chps':
            target_display = '—'
        else:
            target_display = f"{target}{meta['unit']}" if target is not None else '—'

        rows.append({
            'key':    key,
            'label':  meta['label'],
            'target': target_display,
            'type':   row_type,
            'months': [make_kpi_cell(kc['metrics']) for kc in kpi_columns],
        })

    return render(request, 'dashboard/kpi_scorecard.html', {
        'rows':              rows,
        'kpi_columns':       kpi_columns,
        'all_reports':       all_reports,
        'available_months':  available_months,
        'selected_months':   selected_months,
        'selected_county':   selected_county,
        'selected_subcounty': selected_subcounty,
        'counties':          counties,
        'subcounties':       subcounties,
        'has_data':          bool(kpi_columns and report_ids),
        'is_uploader': is_uploader(request.user) if request.user.is_authenticated else False,
    })


# ===========================================================================
# MOH DATA REVIEW
# ===========================================================================

MOH_INDICATORS = [
    # (key, label, section, target, unit, higher_is_better, display_type)
    ('section_workforce', '👥 Workforce',        'header',     None, '', True,  'header'),
    ('active_chps',       'Active CHPs',          'workforce',  None, '',  True, 'simple'),
    ('hh_coverage_pct',   'HH Coverage %',        'workforce',  85,   '%', True, 'simple'),

    ('section_child',     '👶 Child Health',      'header',     None, '', True,  'header'),
    ('u5_assessment_pct', 'U5 Children Assessed', 'child',      100,  '%', True, 'u5'),
    ('total_positive_diag','Positive Diagnoses',  'child',      None, '',  True, 'pos_diag'),
    ('fever_cases',       'Fever',                'child',      None, '',  False,'fever'),

    ('section_iz',        '💉 Immunization',      'header',     None, '', True,  'header'),
    ('iz_defaulters',     'IZ Defaulters',        'iz',         None, '',  False,'simple'),
    ('iz_followup_pct',   'IZ Defaulters Followed Up %','iz',  80,   '%', True, 'simple'),

    ('section_nutrition', '🥗 Nutrition',         'header',     None, '', True,  'header'),
    ('mam_sam_total',     'MAM/SAM Cases',        'nutrition',  None, '',  False,'simple'),
    ('mam_sam_referred_pct','MAM/SAM Referred %', 'nutrition',  90,   '%', True, 'simple'),

    ('section_maternal',  '🤱 Maternal Health',   'header',     None, '', True,  'header'),
    ('pnc_48hr_pct',      'PNC 48hr On-time',     'maternal',   85,   '%', True, 'pnc48'),
    ('pnc_3_7d_pct',      'PNC 3-7d On-time',     'maternal',   85,   '%', True, 'pnc37'),
    ('preg_per_chp',      'Pregnancies Reg/CHP',  'maternal',   1,    '',  True, 'simple'),
    ('facility_delivery_pct','Facility Delivery %','maternal',  80,   '%', True, 'simple'),

    ('section_fp',        '💊 Family Planning',   'header',     None, '', True,  'header'),
    ('fp_wra_assessed',   'WRA Assessed',         'fp',         None, '',  True, 'fp_wra'),
    ('fp_new_users',      'FP New Users',         'fp',         None, '',  True, 'fp_new'),
    ('fp_current_users',  'FP Current Users',     'fp',         None, '',  True, 'simple'),

    ('section_supervision','👀 Supervision',      'header',     None, '', True,  'header'),
    ('supervision_pct',   '% CHPs Supervised',    'supervision',65,   '%', True, 'simple'),

    ('section_art',       '💊 ART',               'header',     None, '', True,  'header'),
    ('art_defaulters',    'ART Defaulters',       'art',        None, '',  False,'simple'),
    ('art_traced',        'ART Traced & Referred','art',        None, '',  False,'simple'),
]
DEFAULT_MOH_TARGETS = {
    'hh_coverage_pct':       85,
    'u5_assessment_pct':     100,
    'iz_followup_pct':       80,
    'mam_sam_referred_pct':  90,
    'pnc_48hr_pct':          85,
    'pnc_3_7d_pct':          85,
    'preg_per_chp':          1,
    'facility_delivery_pct': 80,
    'supervision_pct':       65,
}

QUARTER_MONTHS = {
    'Q1': [1, 2, 3],
    'Q2': [4, 5, 6],
    'Q3': [7, 8, 9],
    'Q4': [10, 11, 12],
}

MONTH_NAMES_FULL = {
    1:'January',2:'February',3:'March',4:'April',5:'May',6:'June',
    7:'July',8:'August',9:'September',10:'October',11:'November',12:'December'
}


def compute_moh_metrics(chw_qs):
    """Compute MOH indicator values from a CHWRecord queryset."""
    if chw_qs is None:
        return None

    active_qs = chw_qs.filter(is_active=True)
    total_all  = chw_qs.count()
    total_active = active_qs.count()

    agg = active_qs.aggregate(
        hh_visits=Sum('hh_visits'),
        reg_hhs=Sum('registered_hhs'),
        u5_assessed=Sum('num_u5_assessed'),
        reg_u5=Sum('registered_children_u5'),
        pos_diag=Sum('positive_diagnoses_u5'),
        fever_cases=Sum('fever_cases'),
        fever_tested=Sum('fever_tested_rdt'),
        iz_def=Sum('iz_defaulters'),
        iz_followup=Sum('iz_defaulters_followed'),
        mam_sam=Sum('mam_sam_total'),
        mam_sam_ref=Sum('mam_sam_referred'),
        pnc_48=Sum('pnc_48hr_ontime'),
        pnc_37=Sum('pnc_3_7d_ontime'),
        deliveries=Sum('total_deliveries'),
        fac_del=Sum('facility_deliveries'),
        preg=Sum('pregnancies_registered'),
        fp_wra=Sum('fp_wra_assessed'),
        fp_new=Sum('fp_new_users'),
        fp_cur=Sum('fp_current_users'),
        sup_visits=Sum('supervision_visits'),
    )

    def pct(num, den):
        if num and den and den > 0:
            return round(num / den * 100, 1)
        return None

    return {
        'active_chps':          total_active,
        'total_chps':           total_all,
        'hh_coverage_pct':      pct(agg['hh_visits'], agg['reg_hhs']),
        # U5 — pct + fraction
        'u5_assessment_pct':    pct(agg['u5_assessed'], agg['reg_u5']),
        'u5_assessed_num':      agg['u5_assessed'] or 0,
        'u5_registered':        agg['reg_u5'] or 0,
        # Positive diagnoses + avg/CHP
        'total_positive_diag':  agg['pos_diag'],
        'pos_diag_per_chp':     round(agg['pos_diag'] / total_active, 2) if agg['pos_diag'] and total_active else None,
        # Fever — cases + tested + tested %
        'fever_cases':          agg['fever_cases'],
        'fever_tested':         agg['fever_tested'],
        'fever_tested_pct':     pct(agg['fever_tested'], agg['fever_cases']),
        # IZ
        'iz_defaulters':        agg['iz_def'],
        'iz_followup_pct':      pct(agg['iz_followup'], agg['iz_def']),
        # Nutrition
        'mam_sam_total':        agg['mam_sam'],
        'mam_sam_referred_pct': pct(agg['mam_sam_ref'], agg['mam_sam']),
        # PNC
        'pnc_48hr_pct':         pct(agg['pnc_48'], agg['deliveries']),
        'pnc_48hr_num':         agg['pnc_48'] or 0,
        'pnc_3_7d_pct':         pct(agg['pnc_37'], agg['deliveries']),
        'pnc_3_7d_num':         agg['pnc_37'] or 0,
        'deliveries':           agg['deliveries'] or 0,
        'preg_per_chp':         round(agg['preg'] / total_active, 2) if agg['preg'] and total_active else None,
        'facility_delivery_pct':pct(agg['fac_del'], agg['deliveries']),
        # FP — counts + percentages
        'fp_wra_assessed':      agg['fp_wra'],
        'fp_wra_pct':           pct(agg['fp_wra'], agg['reg_hhs']),  # % of reg HHs
        'fp_new_users':         agg['fp_new'],
        'fp_new_pct':           pct(agg['fp_new'], agg['fp_wra']),    # % of WRA assessed
        'fp_current_users':     agg['fp_cur'],
        # Supervision
        'supervision_pct':      pct(active_qs.filter(supervised=True).count(), total_active),
        # ART
        'art_defaulters':       None,
        'art_traced':           None,
    }


def aggregate_moh_batches(batches, counties, subcounties, chus):
    """Aggregate CHWRecord data across multiple batches (for quarterly view)."""
    from django.db.models import Sum
    combined = {}

    for b in batches:
        qs = CHWRecord.objects.filter(batch=b)
        if counties:    qs = qs.filter(county__in=counties)
        if subcounties: qs = qs.filter(sub_county__in=subcounties)
        if chus:        qs = qs.filter(community_health_unit__in=chus)
        m = compute_moh_metrics(qs)
        if m is None:
            continue
        for k, v in m.items():
            if v is None:
                continue
            if k not in combined:
                combined[k] = []
            combined[k].append(v)

    if not combined:
        return None

    # For counts: sum; for percentages: average (will be recalculated from raw if possible)
    pct_keys = {'hh_coverage_pct','u5_assessment_pct','iz_followup_pct',
                'mam_sam_referred_pct','pnc_48hr_pct','pnc_3_7d_pct',
                'facility_delivery_pct','supervision_pct'}
    result = {}
    for k, vals in combined.items():
        if k in pct_keys:
            result[k] = round(sum(vals) / len(vals), 1)
        elif k == 'preg_per_chp':
            result[k] = round(sum(vals) / len(vals), 2)
        elif k in ('active_chps', 'total_chps'):
            result[k] = max(vals)  # use latest count not sum
        else:
            result[k] = sum(v for v in vals if isinstance(v, (int, float)))
    return result


@login_required
def moh_review_view(request):
    """MOH Data Review page."""
    view_mode  = request.GET.get('view_mode', 'monthly')  # monthly or quarterly
    counties   = request.GET.getlist('county')
    subcounties= request.GET.getlist('sub_county')
    chus       = request.GET.getlist('chu')

    # Target overrides from POST
    targets = dict(DEFAULT_MOH_TARGETS)
    if request.method == 'POST' and 'save_targets' in request.POST:
        for key in DEFAULT_MOH_TARGETS:
            val = request.POST.get(f'target_{key}', '')
            try:
                targets[key] = float(val) if val else None
            except ValueError:
                pass
        request.session['moh_targets'] = targets
    elif 'moh_targets' in request.session:
        targets.update(request.session['moh_targets'])

    # Get all monthly batches ordered oldest first
    all_monthly = list(UploadBatch.objects.filter(
        period_type='monthly'
    ).order_by('year', 'month'))

    # Build columns
    if view_mode == 'quarterly':
        # Group available monthly batches by quarter/year
        quarter_map = {}
        for b in all_monthly:
            for qname, months in QUARTER_MONTHS.items():
                if b.month in months:
                    key = f"{b.year}-{qname}"
                    if key not in quarter_map:
                        quarter_map[key] = {'label': f"{qname} {b.year}", 'batches': []}
                    quarter_map[key]['batches'].append(b)

        # Override: manually selected quarters
        selected_quarters = request.GET.getlist('quarter')
        if selected_quarters:
            columns = [{'label': quarter_map[q]['label'],
                        'batches': quarter_map[q]['batches'],
                        'key': q}
                       for q in selected_quarters if q in quarter_map]
        else:
            # Default last 4 quarters with data
            all_quarters = sorted(quarter_map.keys())
            columns = [{'label': quarter_map[q]['label'],
                        'batches': quarter_map[q]['batches'],
                        'key': q}
                       for q in all_quarters[-4:]]

        available_periods = [{'key': k, 'label': v['label']} for k, v in sorted(quarter_map.items())]
        selected_periods  = selected_quarters

    else:
        # Monthly: default last 6
        selected_months = request.GET.getlist('batch_month')
        if selected_months:
            batch_map = {str(b.pk): b for b in all_monthly}
            columns = [{'label': batch_map[mid].label, 'batch': batch_map[mid]}
                       for mid in selected_months if mid in batch_map]
        else:
            last6 = all_monthly[-6:]
            columns = [{'label': b.label, 'batch': b} for b in last6]

        available_periods = [{'key': str(b.pk), 'label': b.label} for b in all_monthly]
        selected_periods  = selected_months

    # Compute metrics per column
    col_metrics = []
    for col in columns:
        if view_mode == 'quarterly':
            m = aggregate_moh_batches(col['batches'], counties, subcounties, chus)
        else:
            qs = CHWRecord.objects.filter(batch=col['batch'])
            if counties:    qs = qs.filter(county__in=counties)
            if subcounties: qs = qs.filter(sub_county__in=subcounties)
            if chus:        qs = qs.filter(community_health_unit__in=chus)
            m = compute_moh_metrics(qs)
        col_metrics.append({'label': col['label'], 'metrics': m})

    # Build rows
    def cell_colour(val, target, higher_is_better=True):
        if val is None or target is None:
            return 'grey'
        ratio = val / target * 100 if target else 0
        if higher_is_better:
            if ratio >= 90:  return 'green'
            if ratio >= 50:  return 'yellow'
            return 'red'
        else:
            if ratio <= 110: return 'green'
            if ratio <= 150: return 'yellow'
            return 'red'

    def fmt(val, unit):
        if val is None: return '—'
        if unit == '%': return f"{val}%"
        if isinstance(val, float): return f"{val:,.2f}" if val != int(val) else f"{int(val):,}"
        return f"{val:,}" if isinstance(val, int) else str(val)

    rows = []
    for key, label, section, default_target, unit, hib, display_type in MOH_INDICATORS:
        if display_type == 'header':
            rows.append({'type': 'header', 'label': label, 'key': key})
            continue

        target = targets.get(key, default_target)
        cells  = []
        for cm in col_metrics:
            m   = cm['metrics']
            val = m.get(key) if m else None

            # Build rich display string based on type
            if m is None or val is None:
                display = '—'
            elif display_type == 'u5':
                num = m.get('u5_assessed_num', 0)
                reg = m.get('u5_registered', 0)
                display = f"{val}% ({num:,}/{reg:,})" if reg else f"{val}%"
            elif display_type == 'pos_diag':
                avg = m.get('pos_diag_per_chp')
                display = f"{val:,} ({avg}/CHP)" if avg else f"{val:,}"
            elif display_type == 'fever':
                tested = m.get('fever_tested', 0)
                tpct   = m.get('fever_tested_pct')
                display = f"{val:,} cases, {tested:,} tested ({tpct}%)" if tpct else f"{val:,} cases"
            elif display_type == 'pnc48':
                num = m.get('pnc_48hr_num', 0)
                den = m.get('deliveries', 0)
                display = f"{val}% ({num:,}/{den:,})" if den else f"{val}%"
            elif display_type == 'pnc37':
                num = m.get('pnc_3_7d_num', 0)
                den = m.get('deliveries', 0)
                display = f"{val}% ({num:,}/{den:,})" if den else f"{val}%"
            elif display_type == 'fp_wra':
                wra_pct = m.get('fp_wra_pct')
                display = f"{val:,} ({wra_pct}% of reg. HHs)" if wra_pct else f"{val:,}"
            elif display_type == 'fp_new':
                new_pct = m.get('fp_new_pct')
                display = f"{val:,} ({new_pct}% of WRA)" if new_pct else f"{val:,}"
            elif unit == '%':
                display = f"{val}%"
            elif isinstance(val, float):
                display = f"{val:,.2f}" if val != int(val) else f"{int(val):,}"
            else:
                display = f"{val:,}" if isinstance(val, int) else str(val)

            cells.append({
                'value':   val,
                'display': display,
                'colour':  cell_colour(val, target, hib),
            })

        rows.append({
            'type':   'data',
            'key':    key,
            'label':  label,
            'target': fmt(target, unit) if target is not None else 'TBD',
            'unit':   unit,
            'cells':  cells,
        })

    # Filter options from most recent monthly batch
    filter_opts = {}
    if all_monthly:
        fqs = CHWRecord.objects.filter(batch=all_monthly[-1])
        filter_opts['counties'] = fqs.values_list('county', flat=True).distinct().order_by('county')
        if counties:
            filter_opts['sub_counties'] = fqs.filter(county__in=counties).values_list(
                'sub_county', flat=True).distinct().order_by('sub_county')
        if subcounties:
            filter_opts['chus'] = fqs.filter(sub_county__in=subcounties).values_list(
                'community_health_unit', flat=True).distinct().order_by('community_health_unit')

    return render(request, 'dashboard/moh_review.html', {
        'rows':               rows,
        'col_metrics':        col_metrics,
        'view_mode':          view_mode,
        'available_periods':  available_periods,
        'selected_periods':   selected_periods,
        'filter_opts':        filter_opts,
        'selected_counties':  counties,
        'selected_subcounties': subcounties,
        'selected_chus':      chus,
        'targets':            targets,
        'has_data':           bool(all_monthly),
        'is_uploader': is_uploader(request.user) if request.user.is_authenticated else False,
    })


# ===========================================================================
# PERFORMANCE ACCELERATION SCORECARD
# ===========================================================================

PA_SCORECARD_TARGETS = {
    'cumulative_active_chps': {'label': 'Cumulative Active CHPs',     'target': None, 'unit': '',  'hib': True},
    'weekly_sync_rate':       {'label': 'Weekly Sync Rates/Active CHPs','target':None, 'unit': '',  'hib': True},
    'hh_coverage':            {'label': 'HH Coverage',                'target': 85,   'unit': '%', 'hib': True},
    'child_health':           {'label': 'Child Health Indicators',     'target': None, 'unit': '',  'hib': True},
    'iccm_referrals':         {'label': 'ICCM Referrals',             'target': 90,   'unit': '%', 'hib': True},
    'u2_referrals':           {'label': 'U2 Months Referrals',        'target': 90,   'unit': '%', 'hib': True},
    'iz_completed':           {'label': 'IZ Completed Referrals',     'target': 95,   'unit': '%', 'hib': True},
    'skilled_deliveries':     {'label': 'Skilled Deliveries',         'target': 95,   'unit': '%', 'hib': True},
    'pnc_48hr':               {'label': '48 hrs. ON TIME PNC',        'target': 85,   'unit': '%', 'hib': True},
    'pnc_3_7d':               {'label': '3–7 days PNC',               'target': 85,   'unit': '%', 'hib': True},
    'preg_per_chp':           {'label': 'Preg Reg/CHP',               'target': 1.0,  'unit': '',  'hib': True},
    'preg_visits_vs_active':  {'label': 'Preg Visits vs Active Preg', 'target': 95,   'unit': '%', 'hib': True},
    'wra_assessed':           {'label': 'WRAs (18-49) Assessed',      'target': 90,   'unit': '%', 'hib': True},
    'dash_utilization':       {'label': 'CHAs Dashboard Utilization', 'target': 70,   'unit': '%', 'hib': True},
    'supervision':            {'label': '% CHPs Supervised',          'target': 65,   'unit': '%', 'hib': True},
}


def compute_pa_metrics(chw_qs, sync_qs=None):
    """Compute PA scorecard metrics from CHWRecord queryset."""
    if chw_qs is None:
        return None

    active_qs    = chw_qs.filter(is_active=True)
    total_active = active_qs.count()
    total_all    = chw_qs.count()

    agg = active_qs.aggregate(
        hh_visits    = Sum('hh_visits'),
        reg_hhs      = Sum('registered_hhs'),
        u5_assessed  = Sum('num_u5_assessed'),
        reg_u5       = Sum('registered_children_u5'),
        pos_diag     = Sum('positive_diagnoses_u5'),
        mam_sam      = Sum('mam_sam_total'),
        iccm_ref     = Sum('iccm_referrals_total'),
        iccm_comp    = Sum('iccm_referral_followup'),
        u2_ref       = Sum('iccm_referrals_u2mo'),
        u2_comp      = Sum('iccm_completed_referrals_u2mo'),
        iz_def       = Sum('iz_defaulters'),
        iz_comp      = Sum('iz_defaulters_completed'),
        fac_del      = Sum('facility_deliveries'),
        total_del    = Sum('total_deliveries'),
        pnc_48       = Sum('pnc_48hr_ontime'),
        pnc_37       = Sum('pnc_3_7d_ontime'),
        preg_reg     = Sum('pregnancies_registered'),
        preg_visited = Sum('pregnancies_visited'),
        active_preg  = Sum('active_pregnancies'),
        wra_18_49    = Sum('fp_wra_assessed_18_49'),
        reg_wra      = Sum('registered_hhs'),  # proxy for WRA denominator
    )

    def pct(n, d):
        if n and d and d > 0:
            return round(n / d * 100, 1)
        return None

    # Sync rate from sync_qs if available
    sync_rate = None
    sync_count = None
    if sync_qs is not None:
        synced = sync_qs.filter(days_synced__gte=1).count()
        total_sync = sync_qs.count()
        if total_sync > 0:
            sync_rate  = round(synced / total_sync * 100, 1)
            sync_count = synced

    u5_pct = pct(agg['u5_assessed'], agg['reg_u5'])

    return {
        'cumulative_active_chps': total_active,
        'total_chps':             total_all,
        'weekly_sync_rate':       sync_rate,
        'sync_count':             sync_count,
        'hh_coverage':            pct(agg['hh_visits'], agg['reg_hhs']),
        'hh_visits':              agg['hh_visits'] or 0,
        'reg_hhs':                agg['reg_hhs'] or 0,
        # child health
        'u5_assessed':            agg['u5_assessed'] or 0,
        'reg_u5':                 agg['reg_u5'] or 0,
        'u5_pct':                 u5_pct,
        'pos_diag':               agg['pos_diag'] or 0,
        'pos_diag_per_chp':       round(agg['pos_diag'] / total_active, 2) if agg['pos_diag'] and total_active else None,
        'mam_sam':                agg['mam_sam'] or 0,
        # iCCM
        'iccm_ref_total':         agg['iccm_ref'] or 0,
        'iccm_ref_comp':          agg['iccm_comp'] or 0,
        'iccm_referrals':         pct(agg['iccm_comp'], agg['iccm_ref']),
        # U2
        'u2_ref_total':           agg['u2_ref'] or 0,
        'u2_ref_comp':            agg['u2_comp'] or 0,
        'u2_referrals':           pct(agg['u2_comp'], agg['u2_ref']),
        # IZ
        'iz_defaulters':          agg['iz_def'] or 0,
        'iz_completed':           pct(agg['iz_comp'], agg['iz_def']),
        'iz_comp_num':            agg['iz_comp'] or 0,
        # Deliveries
        'skilled_deliveries':     pct(agg['fac_del'], agg['total_del']),
        'fac_del':                agg['fac_del'] or 0,
        'total_del':              agg['total_del'] or 0,
        # PNC
        'pnc_48hr':               pct(agg['pnc_48'], agg['total_del']),
        'pnc_48_num':             agg['pnc_48'] or 0,
        'pnc_3_7d':               pct(agg['pnc_37'], agg['total_del']),
        'pnc_37_num':             agg['pnc_37'] or 0,
        # Preg
        'preg_per_chp':           round(agg['preg_reg'] / total_active, 2) if agg['preg_reg'] and total_active else None,
        'preg_reg':               agg['preg_reg'] or 0,
        'preg_visits_vs_active':  pct(agg['preg_visited'], agg['active_preg']),
        'preg_visited':           agg['preg_visited'] or 0,
        'active_preg':            agg['active_preg'] or 0,
        # WRA
        'wra_assessed':           pct(agg['wra_18_49'], agg['reg_wra']),
        'wra_18_49':              agg['wra_18_49'] or 0,
        'reg_wra':                agg['reg_wra'] or 0,
        # Supervision
        'supervision':            pct(active_qs.filter(supervised=True).count(), total_active),
        'supervised_count':       active_qs.filter(supervised=True).count(),
        # Dash utilization — filled externally
        'dash_utilization':       None,
    }


def get_dash_util(county, sub_county, report):
    """Get dashboard utilization % for a given geo from a DashUtilReport."""
    if report is None:
        return None
    from .models import DashUtilDataPoint
    dp = DashUtilDataPoint.objects.filter(
        report=report,
        county__iexact=county if county else '',
        sub_county__iexact=sub_county if sub_county else '',
    ).first()
    return dp.utilization_pct if dp else None


def make_pa_cell(metrics, key, target, hib=True, util_pct=None):
    """Build a cell dict for the PA scorecard."""
    if metrics is None:
        return {'display': '—', 'colour': 'grey', 'pct_target': None}

    def colour(val, tgt, hib=True):
        if val is None or tgt is None:
            return 'grey'
        r = val / tgt * 100 if tgt else 0
        if hib:
            return 'green' if r >= 90 else 'yellow' if r >= 50 else 'red'
        else:
            return 'green' if r <= 110 else 'yellow' if r <= 150 else 'red'

    def pct_target(val, tgt):
        if val is None or tgt is None:
            return None
        return round(val / tgt * 100, 1)

    m = metrics

    if key == 'cumulative_active_chps':
        total = m.get('total_chps', 0) or 0
        val   = m.get('cumulative_active_chps', 0) or 0
        pct   = round(val / total * 100, 1) if total else None
        disp  = f"{val:,}/{total:,} ({pct}%)" if pct else f"{val:,}"
        return {'display': disp, 'colour': colour(pct, 100), 'pct_target': pct}

    elif key == 'weekly_sync_rate':
        val  = m.get('weekly_sync_rate')
        cnt  = m.get('sync_count')
        tot  = m.get('cumulative_active_chps', 0) or 0
        disp = f"{cnt:,} ({val}%)" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, 100), 'pct_target': val}

    elif key == 'hh_coverage':
        val  = m.get('hh_coverage')
        vis  = m.get('hh_visits', 0)
        reg  = m.get('reg_hhs', 0)
        disp = f"{val}% ({vis:,})" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    elif key == 'child_health':
        u5p  = m.get('u5_pct')
        au5  = m.get('u5_assessed', 0)
        ru5  = m.get('reg_u5', 0)
        pd_  = m.get('pos_diag', 0)
        pdpc = m.get('pos_diag_per_chp')
        mam  = m.get('mam_sam', 0)
        lines = [
            f"U5 Assessed – {au5:,}/{ru5:,} ({u5p}%)" if u5p else "U5 Assessed – —",
            f"PD – {pdpc}" if pdpc else "PD – —",
            f"MAM/SAM – {mam:,}",
        ]
        return {'display': '\n'.join(lines), 'lines': lines,
                'colour': colour(u5p, 100), 'pct_target': pct_target(u5p, 100)}

    elif key == 'iccm_referrals':
        val  = m.get('iccm_referrals')
        comp = m.get('iccm_ref_comp', 0)
        tot  = m.get('iccm_ref_total', 0)
        disp = f"{comp:,}/{tot:,} ({val}%)" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    elif key == 'u2_referrals':
        val  = m.get('u2_referrals')
        comp = m.get('u2_ref_comp', 0)
        tot  = m.get('u2_ref_total', 0)
        disp = f"{comp:,}/{tot:,} ({val}%)" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    elif key == 'iz_completed':
        val  = m.get('iz_completed')
        comp = m.get('iz_comp_num', 0)
        tot  = m.get('iz_defaulters', 0)
        disp = f"{comp:,}/{tot:,} ({val}%)" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    elif key == 'skilled_deliveries':
        val  = m.get('skilled_deliveries')
        fd   = m.get('fac_del', 0)
        td   = m.get('total_del', 0)
        disp = f"{fd:,}/{td:,} ({val}%)" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    elif key == 'pnc_48hr':
        val  = m.get('pnc_48hr')
        num  = m.get('pnc_48_num', 0)
        td   = m.get('total_del', 0)
        disp = f"{num:,}/{td:,} ({val}%)" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    elif key == 'pnc_3_7d':
        val  = m.get('pnc_3_7d')
        num  = m.get('pnc_37_num', 0)
        td   = m.get('total_del', 0)
        disp = f"{num:,}/{td:,} ({val}%)" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    elif key == 'preg_per_chp':
        val  = m.get('preg_per_chp')
        num  = m.get('preg_reg', 0)
        disp = f"{num:,} ({val})" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    elif key == 'preg_visits_vs_active':
        val  = m.get('preg_visits_vs_active')
        vis  = m.get('preg_visited', 0)
        act  = m.get('active_preg', 0)
        disp = f"{vis:,}/{act:,} ({val}%)" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    elif key == 'wra_assessed':
        val  = m.get('wra_assessed')
        num  = m.get('wra_18_49', 0)
        reg  = m.get('reg_wra', 0)
        disp = f"{num:,}/{reg:,} ({val}%)" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    elif key == 'dash_utilization':
        val  = util_pct
        disp = f"{val}%" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    elif key == 'supervision':
        val  = m.get('supervision')
        cnt  = m.get('supervised_count', 0)
        tot  = m.get('cumulative_active_chps', 0)
        disp = f"{cnt:,}/{tot:,} ({val}%)" if val is not None else '—'
        return {'display': disp, 'colour': colour(val, target), 'pct_target': pct_target(val, target)}

    return {'display': '—', 'colour': 'grey', 'pct_target': None}


@login_required
def pa_scorecard_view(request):
    """Performance Acceleration Scorecard."""
    from .models import DashUtilReport

    selected_counties    = request.GET.getlist('county')
    selected_subcounties = request.GET.getlist('sub_county')
    selected_chus        = request.GET.getlist('chu')

    override_prev_month = request.GET.get('batch_prev_month', '')
    override_weeks      = request.GET.getlist('batch_week')

    all_batches  = UploadBatch.objects.all().order_by('-year', '-month', '-week_start_date')
    all_monthly  = UploadBatch.objects.filter(period_type='monthly').order_by('-year', '-month')

    # Auto-detect
    auto = auto_detect_batches()
    batch_prev_month = UploadBatch.objects.filter(pk=override_prev_month).first() if override_prev_month else auto['prev_month']

    if override_weeks:
        batch_weeks = [UploadBatch.objects.filter(pk=wid).first() for wid in override_weeks if wid]
        batch_weeks = [b for b in batch_weeks if b]
    else:
        batch_weeks = auto['weeks']

    # Dashboard utilization reports
    util_reports = list(DashUtilReport.objects.all().order_by('-year', '-month', '-week'))

    # Find matching util report for each batch
    def find_util_report(batch):
        if not util_reports:
            return None
        if batch.period_type == 'monthly':
            for r in util_reports:
                if r.period_type == 'monthly' and r.year == batch.year and r.month == batch.month:
                    return r
        else:
            if batch.week_start_date:
                for r in util_reports:
                    if r.period_type == 'weekly' and r.week_start_date and abs((r.week_start_date - batch.week_start_date).days) <= 7:
                        return r
        return util_reports[0] if util_reports else None

    def get_chw_qs(batch):
        if batch is None:
            return None
        qs = CHWRecord.objects.filter(batch=batch)
        if selected_counties:    qs = qs.filter(county__in=selected_counties)
        if selected_subcounties: qs = qs.filter(sub_county__in=selected_subcounties)
        if selected_chus:        qs = qs.filter(community_health_unit__in=selected_chus)
        return qs

    def get_sync_qs(chw_batch):
        sync_batch = find_matching_sync_batch(chw_batch)
        if not sync_batch:
            return None
        qs = CHPSyncRecord.objects.filter(batch=sync_batch)
        if selected_counties:    qs = qs.filter(county__in=selected_counties)
        if selected_subcounties: qs = qs.filter(sub_county__in=selected_subcounties)
        if selected_chus:        qs = qs.filter(community_health_unit__in=selected_chus)
        return qs

    # Determine geo for util lookup
    util_county    = selected_counties[0]    if len(selected_counties) == 1    else ''
    util_subcounty = selected_subcounties[0] if len(selected_subcounties) == 1 else ''

    # Compute previous month column
    metrics_prev_month = compute_pa_metrics(get_chw_qs(batch_prev_month), get_sync_qs(batch_prev_month)) if batch_prev_month else None
    util_prev          = find_util_report(batch_prev_month) if batch_prev_month else None
    util_pct_prev      = get_dash_util(util_county, util_subcounty, util_prev)
    if metrics_prev_month:
        metrics_prev_month['dash_utilization'] = util_pct_prev

    # Compute weekly columns
    week_columns = []
    for b in batch_weeks:
        m = compute_pa_metrics(get_chw_qs(b), get_sync_qs(b))
        ur = find_util_report(b)
        up = get_dash_util(util_county, util_subcounty, ur)
        if m:
            m['dash_utilization'] = up
        week_columns.append({'batch': b, 'metrics': m})

    metrics_current = week_columns[-1]['metrics'] if week_columns else None

    # Build rows
    rows = []
    for key, meta in PA_SCORECARD_TARGETS.items():
        target = meta['target']
        prev_cell = make_pa_cell(metrics_prev_month, key, target, meta['hib'])
        week_cells = [make_pa_cell(wc['metrics'], key, target, meta['hib']) for wc in week_columns]

        # % Monthly target achieved from most recent week
        last_cell = week_cells[-1] if week_cells else {'pct_target': None}
        pt = last_cell.get('pct_target')
        pct_colour = 'green' if pt and pt >= 90 else 'yellow' if pt and pt >= 50 else 'red' if pt else 'grey'

        # Dynamic target display matching PPTX format
        ref_metrics = metrics_current or metrics_prev_month or {}
        if key == 'cumulative_active_chps':
            total = ref_metrics.get('total_chps', 0) or 0
            target_display = str(total) if total else 'TBD'
        elif key == 'weekly_sync_rate':
            total = ref_metrics.get('total_chps', 0) or 0
            target_display = str(total) if total else 'TBD'
        elif key == 'hh_coverage':
            reg_hhs = ref_metrics.get('reg_hhs', 0) or 0
            target_display = f"85% ({reg_hhs:,})" if reg_hhs else '85%'
        elif key == 'child_health':
            ru5 = ref_metrics.get('reg_u5', 0) or 0
            target_display = f"U5 Assessed – {ru5:,}\nPD – 10\nMAM/SAM" if ru5 else 'TBD'
        elif key == 'wra_assessed':
            reg_hhs = ref_metrics.get('reg_hhs', 0) or 0
            target_display = f"90% ({reg_hhs:,})" if reg_hhs else '90%'
        else:
            target_display = f"{target}{meta['unit']}" if target is not None else 'TBD'

        rows.append({
            'key':        key,
            'label':      meta['label'],
            'target':     target_display,
            'prev_month': prev_cell,
            'weeks':      week_cells,
            'pct_target': pt,
            'pct_colour': pct_colour,
        })

    # Filter options
    filter_opts  = {}
    filter_batch = batch_weeks[-1] if batch_weeks else batch_prev_month
    if filter_batch:
        fqs = CHWRecord.objects.filter(batch=filter_batch)
        filter_opts['counties'] = fqs.values_list('county', flat=True).distinct().order_by('county')
        if selected_counties:
            filter_opts['sub_counties'] = fqs.filter(county__in=selected_counties).values_list('sub_county', flat=True).distinct().order_by('sub_county')
        if selected_subcounties:
            filter_opts['chus'] = fqs.filter(sub_county__in=selected_subcounties).values_list('community_health_unit', flat=True).distinct().order_by('community_health_unit')

    return render(request, 'dashboard/pa_scorecard.html', {
        'rows':               rows,
        'batch_prev_month':   batch_prev_month,
        'week_batches':       [wc['batch'] for wc in week_columns],
        'all_batches':        all_batches,
        'all_monthly':        all_monthly,
        'util_reports':       util_reports,
        'override_prev_month': override_prev_month,
        'override_weeks':     override_weeks,
        'filter_opts':        filter_opts,
        'selected_counties':  selected_counties,
        'selected_subcounties': selected_subcounties,
        'selected_chus':      selected_chus,
        'has_data':           bool(batch_weeks or batch_prev_month),
        'is_uploader': is_uploader(request.user) if request.user.is_authenticated else False,
    })


@login_required
def dash_util_upload_view(request):
    """Upload Dashboard Utilization Report."""
    from .models import DashUtilReport
    from .dash_util_parser import parse_dash_util_report

    reports = DashUtilReport.objects.all()
    success = error = None

    if request.method == 'POST':
        f           = request.FILES.get('file')
        period_type = request.POST.get('period_type', 'weekly')
        period_label = request.POST.get('period_label', '')
        year        = request.POST.get('year', '')
        month       = request.POST.get('month', '') or None
        week        = request.POST.get('week', '') or None
        week_start  = request.POST.get('week_start_date', '') or None

        if f and period_label and year:
            from datetime import date as date_cls
            report = DashUtilReport(
                file=f, period_type=period_type, period_label=period_label,
                year=int(year),
                month=int(month) if month else None,
                week=int(week) if week else None,
                week_start_date=date_cls.fromisoformat(week_start) if week_start else None,
                uploaded_by=request.user,
            )
            report.save()
            rows, errors = parse_dash_util_report(report)
            if errors:
                error = f"Uploaded with warnings: {'; '.join(errors[:3])}"
            else:
                success = f"Dashboard Utilization Report uploaded — {rows} data points parsed."
        else:
            error = "Please fill all required fields."

    return render(request, 'dashboard/dash_util_upload.html', {
        'reports': reports, 'success': success, 'error': error,
        'is_uploader': is_uploader(request.user),
    })


@login_required
def dash_util_delete_view(request, pk):
    from .models import DashUtilReport
    if request.method == 'POST':
        get_object_or_404(DashUtilReport, pk=pk).delete()
    return redirect('dash_util_upload')



@login_required
def pa_scorecard_pptx(request):
    """Generate and download PA scorecard as PPTX using python-pptx."""
    import io, re
    from django.http import HttpResponse
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def rgb(h):
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    NAVY='1B3A6B'; WHITE='FFFFFF'
    GREEN='00B050'; GREEN_T='FFFFFF'
    YELLOW='FFFF00'; YELLOW_T='000000'
    RED='FF0000'; RED_T='FFFFFF'
    GREY='D3D3D3'; GREY_T='000000'
    GREY_H='D3D3D3'; GREY_TGT='D3D3D3'
    PCT_GRN='00B050'; PCT_GRN_T='FFFFFF'
    PCT_YLW='FFFF00'; PCT_YLW_T='000000'
    PCT_RED='FF0000'; PCT_RED_T='FFFFFF'

    def cell_bg_txt(colour):
        return {'green':(GREEN,GREEN_T),'yellow':(YELLOW,YELLOW_T),'red':(RED,RED_T)}.get(colour,(GREY,GREY_T))

    def pct_bg_txt(colour):
        return {'green':(PCT_GRN,PCT_GRN_T),'yellow':(PCT_YLW,PCT_YLW_T),'red':(PCT_RED,PCT_RED_T)}.get(colour,(GREY,GREY_T))

    def shorten(label):
        if not label: return ''
        label = str(label)
        if 'Monthly' in label or 'monthly' in label:
            m = re.search(r'([A-Za-z]+)\s+\d{4}', label)
            return m.group(1) if m else label[:8]
        m = re.search(r'(\d+)\s+([A-Za-z]+)\s+\d{4}', label)
        return (m.group(2)[:3] + ' WK') if m else label[:10]

    # Recompute rows
    selected_counties    = request.GET.getlist('county')
    selected_subcounties = request.GET.getlist('sub_county')
    selected_chus        = request.GET.getlist('chu')
    override_prev_month  = request.GET.get('batch_prev_month', '')
    override_weeks       = request.GET.getlist('batch_week')

    from .models import DashUtilReport
    auto = auto_detect_batches()
    batch_prev_month = UploadBatch.objects.filter(pk=override_prev_month).first() if override_prev_month else auto['prev_month']
    batch_weeks = [UploadBatch.objects.filter(pk=wid).first() for wid in override_weeks if wid] if override_weeks else auto['weeks']
    batch_weeks = [b for b in batch_weeks if b]

    util_reports   = list(DashUtilReport.objects.all().order_by('-year','-month','-week'))
    util_county    = selected_counties[0]    if len(selected_counties) == 1    else ''
    util_subcounty = selected_subcounties[0] if len(selected_subcounties) == 1 else ''

    def get_chw_qs(batch):
        if not batch: return None
        qs = CHWRecord.objects.filter(batch=batch)
        if selected_counties:    qs = qs.filter(county__in=selected_counties)
        if selected_subcounties: qs = qs.filter(sub_county__in=selected_subcounties)
        if selected_chus:        qs = qs.filter(community_health_unit__in=selected_chus)
        return qs

    def get_sync_qs(b):
        sb = find_matching_sync_batch(b)
        if not sb: return None
        qs = CHPSyncRecord.objects.filter(batch=sb)
        if selected_counties:    qs = qs.filter(county__in=selected_counties)
        if selected_subcounties: qs = qs.filter(sub_county__in=selected_subcounties)
        return qs

    def find_util(batch):
        if not util_reports or not batch: return None
        if batch.period_type == 'monthly':
            for r in util_reports:
                if r.period_type == 'monthly' and r.year == batch.year and r.month == batch.month: return r
        elif batch.week_start_date:
            for r in util_reports:
                if r.period_type == 'weekly' and r.week_start_date and abs((r.week_start_date - batch.week_start_date).days) <= 7: return r
        return util_reports[0] if util_reports else None

    metrics_pm = compute_pa_metrics(get_chw_qs(batch_prev_month), get_sync_qs(batch_prev_month)) if batch_prev_month else None
    if metrics_pm: metrics_pm['dash_utilization'] = get_dash_util(util_county, util_subcounty, find_util(batch_prev_month))

    week_cols = []
    for b in batch_weeks:
        m = compute_pa_metrics(get_chw_qs(b), get_sync_qs(b))
        if m: m['dash_utilization'] = get_dash_util(util_county, util_subcounty, find_util(b))
        week_cols.append({'batch': b, 'metrics': m})

    rows = []
    ref_metrics = (week_cols[-1]['metrics'] if week_cols else None) or metrics_pm or {}
    for key, meta in PA_SCORECARD_TARGETS.items():
        target = meta['target']
        pm_cell  = make_pa_cell(metrics_pm, key, target, meta['hib'])
        wk_cells = [make_pa_cell(wc['metrics'], key, target, meta['hib']) for wc in week_cols]
        last = wk_cells[-1] if wk_cells else {}
        pt = last.get('pct_target')
        pc = 'green' if pt and pt >= 90 else 'yellow' if pt and pt >= 50 else 'red' if pt else 'grey'

        # Dynamic target display matching web view
        if key == 'cumulative_active_chps':
            total = ref_metrics.get('total_chps', 0) or 0
            target_display = str(total) if total else 'TBD'
        elif key == 'weekly_sync_rate':
            total = ref_metrics.get('total_chps', 0) or 0
            target_display = str(total) if total else 'TBD'
        elif key == 'hh_coverage':
            reg_hhs = ref_metrics.get('reg_hhs', 0) or 0
            target_display = f"85% ({reg_hhs:,})" if reg_hhs else '85%'
        elif key == 'child_health':
            ru5 = ref_metrics.get('reg_u5', 0) or 0
            target_display = f"U5 Assessed – {ru5:,}\nPD – 10\nMAM/SAM" if ru5 else 'TBD'
        elif key == 'wra_assessed':
            reg_hhs = ref_metrics.get('reg_hhs', 0) or 0
            target_display = f"90% ({reg_hhs:,})" if reg_hhs else '90%'
        else:
            target_display = f"{target}{meta['unit']}" if target is not None else 'TBD'

        rows.append({'key':key,'label':meta['label'],'target':target_display,
                     'prev_month':pm_cell,'weeks':wk_cells,'pct_target':pt,'pct_colour':pc})

    # Build PPTX
    prs = Presentation()
    prs.slide_width  = Inches(13.3)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    MARGIN = Inches(0.12)
    W_TOTAL = prs.slide_width  - MARGIN * 2
    H_TOTAL = prs.slide_height - MARGIN * 2
    N_WEEKS = len(week_cols)

    W_IND  = Inches(1.5);  W_TGT = Inches(0.82);  W_PCT = Inches(0.92)
    W_DATA = int((W_TOTAL - W_IND - W_TGT - W_PCT) / (1 + N_WEEKS))
    col_ws = [W_IND, W_TGT] + [W_DATA] * (1 + N_WEEKS) + [W_PCT]

    H_HDR   = Inches(0.42)
    H_AVAIL = H_TOTAL - H_HDR
    child_idx = next((i for i,r in enumerate(rows) if r['key']=='child_health'), -1)
    n_norm    = len(rows) - (1 if child_idx>=0 else 0)
    H_NORM    = int(H_AVAIL / (n_norm + 2.2))
    row_hs    = [int(H_NORM*2.2) if r['key']=='child_health' else H_NORM for r in rows]

    xs = [MARGIN]
    for w in col_ws[:-1]: xs.append(xs[-1]+w)
    ys = [MARGIN+H_HDR]
    for h in row_hs[:-1]: ys.append(ys[-1]+h)

    def add_cell(x, y, w, h, text, bg, txt_color, bold=True,
                 font_size=8, align=PP_ALIGN.CENTER, lines=None):
        shape = slide.shapes.add_shape(1, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(bg)
        shape.line.color.rgb = rgb(WHITE)
        shape.line.width = Pt(1.2)
        tf = shape.text_frame
        tf.word_wrap = True
        if lines:
            for i, line in enumerate(lines):
                p = tf.paragraphs[i] if i==0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = line; run.font.size=Pt(font_size-0.5)
                run.font.color.rgb=rgb(txt_color); run.font.bold=(i==0); run.font.name='Calibri'
        else:
            p = tf.paragraphs[0]; p.alignment = align
            run = p.add_run()
            run.text = str(text) if text else ''
            run.font.size=Pt(font_size); run.font.color.rgb=rgb(txt_color)
            run.font.bold=bold; run.font.name='Calibri'

    # Header
    geo_parts = selected_subcounties or selected_counties or ['Kenya']
    h_texts = ['/'.join(geo_parts)+'\nIndicator','Monthly\nTarget',
               shorten(batch_prev_month.label if batch_prev_month else 'Prev Month')]
    h_texts += [shorten(wc['batch'].label) for wc in week_cols]
    h_texts += ['% of Monthly\nTarget Achieved']
    for ci,(hx,hw,ht) in enumerate(zip(xs,col_ws,h_texts)):
        add_cell(hx,MARGIN,hw,H_HDR,ht,NAVY,WHITE,bold=True,font_size=8,
                 align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)

    # Data rows
    for ri,(row,ry,rh) in enumerate(zip(rows,ys,row_hs)):
        is_child = row['key']=='child_health'
        add_cell(xs[0],ry,col_ws[0],rh,row['label'],GREY_H,NAVY,bold=True,font_size=7.5,align=PP_ALIGN.LEFT)
        add_cell(xs[1],ry,col_ws[1],rh,row['target'],GREY_TGT,'374151',bold=True,font_size=7.5)
        all_cells = [row['prev_month']] + row['weeks']
        for ci,cell in enumerate(all_cells):
            col_i = ci+2
            bg,txt = cell_bg_txt(cell.get('colour','grey'))
            lines  = cell.get('lines') if is_child and cell.get('lines') else None
            disp   = '' if is_child else cell.get('display','')
            add_cell(xs[col_i],ry,col_ws[col_i],rh,disp,bg,txt,bold=True,
                     font_size=6.5 if is_child else 7.5,align=PP_ALIGN.CENTER,lines=lines)
        bg,txt = pct_bg_txt(row['pct_colour'])
        pct_ci = len(col_ws)-1
        add_cell(xs[pct_ci],ry,col_ws[pct_ci],rh,
                 f"{row['pct_target']}%" if row['pct_target'] else '',
                 bg,txt,bold=True,font_size=8)

    buf = io.BytesIO()
    prs.save(buf); buf.seek(0)
    filename = f"PA_Scorecard_{'_'.join(geo_parts)}.pptx".replace(' ','_')
    response = HttpResponse(buf.getvalue(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


@login_required
def admin_reprocess_view(request):
    """Temporary admin endpoint to reprocess all batches on Railway."""
    if not request.user.is_staff:
        from django.http import HttpResponseForbidden
        return HttpResponseForbidden("Staff only.")

    from .parsers import parse_chw_file, parse_supervision_file
    from django.utils.timezone import now

    results = []
    batches = UploadBatch.objects.all().order_by('year', 'month', 'week_start_date')

    for batch in batches:
        batch_results = {'label': batch.label, 'chw': None, 'sup': None, 'errors': []}

        # Reprocess CHW file
        if batch.chw_file:
            try:
                CHWRecord.objects.filter(batch=batch).delete()
                batch.chw_file.seek(0)
                rows, errors = parse_chw_file(batch, batch.chw_file)
                batch_results['chw'] = f"{rows} rows"
                if errors:
                    batch_results['errors'] += errors[:3]
            except Exception as e:
                batch_results['errors'].append(f"CHW error: {e}")

        # Reprocess supervision file
        if batch.supervision_file:
            try:
                from .models import SupervisionRecord
                SupervisionRecord.objects.filter(batch=batch).delete()
                batch.supervision_file.seek(0)
                rows, errors = parse_supervision_file(batch, batch.supervision_file)
                batch_results['sup'] = f"{rows} rows"
                if errors:
                    batch_results['errors'] += errors[:3]
            except Exception as e:
                batch_results['errors'].append(f"Supervision error: {e}")

        # Recompute indicators
        try:
            from .parsers import compute_indicators
            compute_indicators(batch)
        except Exception as e:
            batch_results['errors'].append(f"Indicators error: {e}")

        results.append(batch_results)

    # Render simple HTML result
    html = '<html><body style="font-family:monospace;padding:20px;">'
    html += f'<h2>Reprocess Results — {now().strftime("%Y-%m-%d %H:%M:%S")}</h2>'
    html += f'<p>Processed {len(results)} batches</p><table border="1" cellpadding="6">'
    html += '<tr><th>Batch</th><th>CHW Rows</th><th>Supervision Rows</th><th>Errors</th></tr>'
    for r in results:
        err_str = '<br>'.join(r['errors']) if r['errors'] else '✅ OK'
        html += f"<tr><td>{r['label']}</td><td>{r['chw'] or '—'}</td><td>{r['sup'] or '—'}</td><td>{err_str}</td></tr>"
    html += '</table><br><a href="/">← Back to Dashboard</a></body></html>'
    return HttpResponse(html)